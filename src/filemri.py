#!/usr/bin/env python3
"""
DeepFileX - Advanced File Analysis System
차세대 파일 검색 및 분석 도구
"""

import sys
import os
import sqlite3
import logging
import threading
import time
import json
import pickle
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
import hashlib

# Configure logging with proper path handling
import tempfile
import atexit

# Create logs directory (override with DEEPFILEX_LOG_DIR to avoid permission issues)
log_base = os.environ.get('DEEPFILEX_LOG_DIR')
if log_base:
    log_dir = Path(log_base)
else:
    log_dir = Path.home() / 'AppData' / 'Roaming' / 'DeepFileX'
log_dir.mkdir(parents=True, exist_ok=True)
log_file = log_dir / 'deepfilex.log'

# Configure logging handlers with UTF-8 encoding for emoji support
stream_handler = logging.StreamHandler(sys.stdout)
stream_handler.setLevel(logging.INFO)
stream_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))

# Set UTF-8 encoding for the stream handler to support emoji
try:
    import io
    if hasattr(sys.stdout, 'buffer'):
        stream_handler.stream = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
except Exception:
    # Fallback: use errors='replace' to avoid crashes
    stream_handler.stream = sys.stdout

file_handler = logging.FileHandler(str(log_file), encoding='utf-8')
file_handler.setLevel(logging.INFO)
file_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[stream_handler, file_handler]
)
logger = logging.getLogger(__name__)

# PyQt6 imports
try:
    from PyQt6.QtWidgets import (
        QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
        QPushButton, QLineEdit, QListWidget, QListWidgetItem, QTextEdit,
        QLabel, QFileDialog, QMessageBox, QSplitter, QProgressBar,
        QTabWidget, QComboBox, QCheckBox, QSpinBox, QGroupBox,
        QPlainTextEdit, QDialog, QDialogButtonBox, QFormLayout
    )
    from PyQt6.QtCore import Qt, QThread, pyqtSignal, QTimer, QSettings
    from PyQt6.QtGui import QFont, QIcon
    PYQT_AVAILABLE = True
except ImportError:
    PYQT_AVAILABLE = False

# 🆕 SmartLinks 수익화 시스템 통합 (Adsterra)
try:
    from filemri_smartlinks import SmartLinksAdWidget as AdBanner
    from filemri_smartlinks import DeepFileXSmartLinksManager
    SMARTLINKS_AVAILABLE = True
    logger.info("Adsterra SmartLinks system loaded")
except ImportError as e:
    SMARTLINKS_AVAILABLE = False
    logger.warning(f"SmartLinks module not available: {e}")

# 🆕 자동 업데이트 시스템 통합
try:
    from update_checker import UpdateChecker, UpdateDialog
    from version_config import VERSION, UPDATE_CONFIG, UPDATE_MESSAGES
    UPDATE_SYSTEM_AVAILABLE = True
    logger.info("✅ 업데이트 시스템 로드 성공")
except ImportError as e:
    UPDATE_SYSTEM_AVAILABLE = False
    logger.warning(f"⚠️ 업데이트 시스템 모듈 없음: {e}")

# Import for file content extraction
try:
    import chardet
    CHARDET_AVAILABLE = True
except ImportError:
    CHARDET_AVAILABLE = False

# PDF content extraction
try:
    import PyPDF2
    PDF_AVAILABLE = True
    PDF_METHOD = 'PyPDF2'
    logger.info("PDF support: PyPDF2 available")
except ImportError:
    try:
        import pdfplumber
        PDF_AVAILABLE = True
        PDF_METHOD = 'pdfplumber'
        logger.info("PDF support: pdfplumber available")
    except ImportError:
        PDF_AVAILABLE = False
        PDF_METHOD = None
        logger.warning("PDF support: No PDF library available")

# Office document extraction
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import zipfile
    ZIP_AVAILABLE = True
except ImportError:
    ZIP_AVAILABLE = False

class TurboFileExtractor:
    """Ultra-fast file content extractor with PDF support"""
    
    def __init__(self):
        self.supported_extensions = {
            '.txt', '.md', '.log', '.csv', '.json', '.xml', '.html', '.css',
            '.js', '.py', '.java', '.c', '.cpp', '.h', '.cs', '.php', '.rb',
            '.go', '.rs', '.swift', '.kt', '.sql', '.sh', '.bat', '.ps1',
            '.yaml', '.yml', '.ini', '.cfg', '.conf', '.toml', '.pdf',
            # Office documents
            '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
            # Email and other formats
            '.pst', '.eml', '.msg',
            # Korean documents
            '.hwp', '.hwpx',
            # Archive files (filename search only)
            '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2',
            # Image files (filename search only)
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg',
            # Other text-based formats
            '.rtf', '.tex', '.scss', '.less', '.vue', '.tsx', '.jsx'
        }
        # Small file cache for speed
        self._cache = {}
        self._cache_max_size = 500
    
    def can_extract_fast(self, file_path: str) -> bool:
        """Lightning fast file check with generous size limits"""
        try:
            path = Path(file_path)
            if not path.exists() or not path.is_file():
                return False
            
            ext = path.suffix.lower()
            
            # More generous size limits for different file types
            if ext == '.pdf':
                size_limit = 50 * 1024 * 1024  # 50MB for PDF
            elif ext in ['.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls']:
                size_limit = 100 * 1024 * 1024  # 100MB for Office docs
            elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                size_limit = 500 * 1024 * 1024  # 500MB for archives
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                size_limit = 20 * 1024 * 1024  # 20MB for images
            elif ext in ['.hwp', '.hwpx']:
                size_limit = 50 * 1024 * 1024  # 50MB for Korean docs
            else:
                size_limit = 10 * 1024 * 1024  # 10MB for text files
            
            if path.stat().st_size > size_limit:
                logger.info(f"File too large, skipping: {file_path} ({path.stat().st_size / 1024 / 1024:.1f}MB)")
                return False
            
            return ext in self.supported_extensions
        except Exception as e:
            logger.warning(f"Error checking file {file_path}: {e}")
            return False
    
    def extract_fast(self, file_path: str) -> Dict[str, Any]:
        """Ultra-fast content extraction with comprehensive file support"""
        try:
            path = Path(file_path)
            if not self.can_extract_fast(file_path):
                logger.debug(f"Skipping unsupported file: {file_path}")
                return {'success': False, 'content': '', 'metadata': {}}
            
            stat = path.stat()
            
            # Cache check
            cache_key = f"{file_path}_{stat.st_mtime}"
            if cache_key in self._cache:
                return self._cache[cache_key]
            
            metadata = {
                'filename': path.name,
                'extension': path.suffix.lower(),
                'size': stat.st_size,
                'modified': datetime.fromtimestamp(stat.st_mtime),
                'created': datetime.fromtimestamp(stat.st_ctime)
            }
            
            # Extract content based on file type
            ext = path.suffix.lower()
            logger.debug(f"Processing {ext} file: {path.name}")
            
            if ext == '.pdf':
                content = self._extract_pdf_content(file_path)
            elif ext in ['.docx', '.doc']:
                content = self._extract_word_content(file_path)
            elif ext in ['.pptx', '.ppt']:
                content = self._extract_pptx_content(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = self._extract_excel_content(file_path)
            elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                content = self._extract_archive_filelist(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg']:
                # Image files - create meaningful content for filename search with extra keywords
                base_name = path.stem.lower().replace('_', ' ').replace('-', ' ')
                img_keywords = []
                
                # Add common image-related keywords based on filename patterns
                if any(word in base_name for word in ['photo', 'pic', 'image', 'img']):
                    img_keywords.append('photo')
                if any(word in base_name for word in ['screen', 'shot', 'capture']):
                    img_keywords.append('screenshot')
                if any(word in base_name for word in ['logo', 'icon', 'symbol']):
                    img_keywords.append('graphic')
                if any(word in base_name for word in ['wall', 'paper', 'background']):
                    img_keywords.append('wallpaper')
                
                extra_keywords = ' '.join(img_keywords)
                content = f"Image file: {path.name} {base_name} {extra_keywords} picture visual graphic"
                logger.debug(f"Processing image file: {path.name} with keywords: {extra_keywords}")
            elif ext in ['.pst', '.eml', '.msg', '.hwp', '.hwpx']:
                content = f"Document file: {path.name} {path.stem.lower().replace('_', ' ').replace('-', ' ')}"  # For filename search only
            else:
                content = self._read_text_ultra_fast(file_path)
            
            # Ensure we have some content for indexing (even if just filename)
            if not content or content.strip() == "":
                content = f"File: {path.name}"
            
            result = {
                'success': True,
                'content': content[:10000],  # Limit for speed
                'metadata': metadata
            }
            
            logger.debug(f"Successfully extracted content from {path.name}: {len(content)} chars")
            
            # Cache management
            if len(self._cache) >= self._cache_max_size:
                # Remove oldest entry
                oldest_key = next(iter(self._cache))
                del self._cache[oldest_key]
            
            self._cache[cache_key] = result
            return result
            
        except Exception as e:
            logger.warning(f"Fast extraction failed for {file_path}: {e}")
            # Still return success with filename for filename search
            try:
                return {
                    'success': True,
                    'content': f"File: {Path(file_path).name}",
                    'metadata': {
                        'filename': Path(file_path).name,
                        'extension': Path(file_path).suffix.lower(),
                        'size': 0,
                        'modified': datetime.now(),
                        'created': datetime.now()
                    }
                }
            except:
                return {'success': False, 'content': '', 'metadata': {}}
    
    def _extract_word_content(self, file_path: str) -> str:
        """Extract text from Word documents"""
        if not DOCX_AVAILABLE:
            return f"Word document: {Path(file_path).name}"
        
        try:
            if file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                text = ""
                for paragraph in doc.paragraphs[:20]:  # First 20 paragraphs for speed
                    text += paragraph.text + "\n"
                return text
            else:
                # .doc files need python-docx2txt or similar
                return f"Word document: {Path(file_path).name}"
        except Exception as e:
            logger.warning(f"Word extraction failed for {file_path}: {e}")
            return f"Word document: {Path(file_path).name}"
    
    def _extract_excel_content(self, file_path: str) -> str:
        """Extract text from Excel files"""
        if not EXCEL_AVAILABLE:
            return f"Excel document: {Path(file_path).name}"
        
        try:
            if file_path.endswith('.xlsx'):
                workbook = openpyxl.load_workbook(file_path, read_only=True)
                text = ""
                # Get first few sheets
                for sheet_name in list(workbook.sheetnames)[:3]:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows(max_row=50, values_only=True):  # First 50 rows
                        for cell in row:
                            if cell and isinstance(cell, str):
                                text += str(cell) + " "
                    text += "\n"
                workbook.close()
                return text
            else:
                # .xls files need xlrd or similar
                return f"Excel document: {Path(file_path).name}"
        except Exception as e:
            logger.warning(f"Excel extraction failed for {file_path}: {e}")
            return f"Excel document: {Path(file_path).name}"
    
    def _extract_pptx_content(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        if not PPTX_AVAILABLE:
            return f"PowerPoint document: {Path(file_path).name}"
        
        try:
            if file_path.endswith('.pptx'):
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides[:20]:  # First 20 slides for speed
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            else:
                # .ppt files need python-pptx or similar
                return f"PowerPoint document: {Path(file_path).name}"
        except Exception as e:
            logger.warning(f"PowerPoint extraction failed for {file_path}: {e}")
            return f"PowerPoint document: {Path(file_path).name}"
    
    def _extract_archive_filelist(self, file_path: str) -> str:
        """Extract file list from archives for filename search"""
        if not ZIP_AVAILABLE:
            return f"Archive: {Path(file_path).name}"
        
        try:
            if file_path.endswith('.zip'):
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    filenames = zip_file.namelist()
                    # Return list of filenames for searching
                    return f"Archive contents: {' '.join(filenames[:100])}"  # First 100 files
            else:
                return f"Archive: {Path(file_path).name}"
        except Exception as e:
            logger.warning(f"Archive extraction failed for {file_path}: {e}")
            return f"Archive: {Path(file_path).name}"
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """Extract text from PDF files"""
        if not PDF_AVAILABLE:
            logger.debug(f"PDF extraction skipped - no library available: {file_path}")
            return f"PDF document: {Path(file_path).name}"
        
        try:
            logger.debug(f"Attempting PDF extraction using {PDF_METHOD}: {file_path}")
            
            if PDF_METHOD == 'pdfplumber':
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page_num, page in enumerate(pdf.pages[:5]):  # First 5 pages for speed
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                        except Exception as e:
                            logger.warning(f"Error extracting page {page_num} from {file_path}: {e}")
                    if text.strip():
                        logger.debug(f"Successfully extracted {len(text)} chars from PDF: {Path(file_path).name}")
                        return text
                    else:
                        logger.warning(f"No text extracted from PDF: {file_path}")
                        return f"PDF document: {Path(file_path).name}"
            else:  # PyPDF2
                import PyPDF2
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page_num in range(min(5, len(reader.pages))):  # First 5 pages
                        try:
                            page = reader.pages[page_num]
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                        except Exception as e:
                            logger.warning(f"Error extracting page {page_num} from {file_path}: {e}")
                    if text.strip():
                        logger.debug(f"Successfully extracted {len(text)} chars from PDF: {Path(file_path).name}")
                        return text
                    else:
                        logger.warning(f"No text extracted from PDF: {file_path}")
                        return f"PDF document: {Path(file_path).name}"
                        
        except Exception as e:
            logger.warning(f"PDF extraction failed for {file_path}: {e}")
            return f"PDF document: {Path(file_path).name}"
    
    def _read_text_ultra_fast(self, file_path: str) -> str:
        """Ultra-fast text reading with encoding detection"""
        try:
            # Try UTF-8 first (fastest)
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(10000)  # Read first 10KB for speed
        except:
            try:
                # Fallback with detection if available
                if CHARDET_AVAILABLE:
                    with open(file_path, 'rb') as f:
                        raw_data = f.read(1024)
                        encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
                    
                    with open(file_path, 'r', encoding=encoding, errors='ignore') as f:
                        return f.read(10000)
                else:
                    # Last resort
                    with open(file_path, 'r', encoding='latin-1') as f:
                        return f.read(10000)
            except:
                return ""


class TurboSearchEngine:
    """Ultra-fast in-memory search engine with persistence and filename search"""
    
    def __init__(self, db_path: str = "turbo_search.db"):
        self.db_path = db_path
        self.index_data = {}  # file_path -> {content, metadata, keywords, filename_keywords}
        self.word_index = {}  # word -> set of file_paths (content words)
        self.filename_index = {}  # word -> set of file_paths (filename words)
        self.stats = {'total_files': 0, 'total_words': 0, 'total_filename_words': 0, 'last_update': None}
        self._lock = threading.RLock()
        
        # Initialize database
        self._init_database()
        
    def _init_database(self):
        """Initialize SQLite database"""
        try:
            with sqlite3.connect(self.db_path) as conn:
                conn.execute('''
                    CREATE TABLE IF NOT EXISTS files (
                        path TEXT PRIMARY KEY,
                        content TEXT,
                        metadata TEXT,
                        keywords TEXT,
                        filename_keywords TEXT,
                        indexed_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                    )
                ''')
                conn.execute('''
                    CREATE INDEX IF NOT EXISTS idx_content ON files(content)
                ''')
                conn.execute('''
                    CREATE INDEX IF NOT EXISTS idx_keywords ON files(keywords)
                ''')
                conn.execute('''
                    CREATE INDEX IF NOT EXISTS idx_filename_keywords ON files(filename_keywords)
                ''')
                conn.commit()
        except Exception as e:
            logger.error(f"Database initialization failed: {e}")
    
    def save_index(self, filepath: str = None) -> bool:
        """Save current index to file"""
        try:
            if not filepath:
                filepath = f"turbo_index_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pkl"
            
            save_data = {
                'index_data': self.index_data,
                'word_index': {k: list(v) for k, v in self.word_index.items()},  # Convert sets to lists
                'filename_index': {k: list(v) for k, v in self.filename_index.items()},  # Convert sets to lists
                'stats': self.stats,
                'version': '1.1',  # Updated version for filename support
                'created': datetime.now().isoformat()
            }
            
            with open(filepath, 'wb') as f:
                pickle.dump(save_data, f, protocol=pickle.HIGHEST_PROTOCOL)
            
            logger.info(f"Index saved to {filepath}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to save index: {e}")
            return False
    
    def load_index(self, filepath: str) -> bool:
        """Load index from file with enhanced error handling and logging"""
        try:
            logger.info(f"TurboSearchEngine.load_index starting for: {filepath}")
            
            # Step 1: Validate file
            if not os.path.exists(filepath):
                logger.error(f"Index file does not exist: {filepath}")
                return False
                
            file_size = os.path.getsize(filepath)
            logger.info(f"Index file size: {file_size:,} bytes")
            
            # Step 2: Load and validate data
            logger.info("Loading pickle data...")
            
            with open(filepath, 'rb') as f:
                save_data = pickle.load(f)
            
            logger.info(f"Pickle data loaded successfully, type: {type(save_data)}")
            
            if not isinstance(save_data, dict):
                logger.error(f"Invalid index data type: {type(save_data)}")
                return False
            
            # Step 3: Validate version
            version = save_data.get('version', '1.0')
            logger.info(f"Index version: {version}")
            
            if version not in ['1.0', '1.1']:
                logger.warning(f"Index version mismatch: {version}")
                return False
            
            # Step 4: Load data with thread safety
            logger.info("Updating search engine data...")
            
            with self._lock:
                # Load index data
                self.index_data = save_data.get('index_data', {})
                logger.info(f"Loaded {len(self.index_data)} index entries")
                
                # Convert lists back to sets for word index
                word_index = save_data.get('word_index', {})
                self.word_index = {k: set(v) for k, v in word_index.items()}
                logger.info(f"Loaded {len(self.word_index)} word index entries")
                
                # Handle filename index (new in v1.1)
                if version == '1.1':
                    filename_index = save_data.get('filename_index', {})
                    self.filename_index = {k: set(v) for k, v in filename_index.items()}
                    logger.info(f"Loaded {len(self.filename_index)} filename index entries")
                else:
                    # Rebuild filename index for old version
                    logger.info("Rebuilding filename index for older version...")
                    self.filename_index = {}
                    self._rebuild_filename_index()
                    logger.info(f"Rebuilt {len(self.filename_index)} filename index entries")
                
                # Load statistics
                self.stats = save_data.get('stats', {
                    'total_files': 0, 'total_words': 0, 'total_filename_words': 0, 'last_update': None
                })
                
                # Add missing stats for old versions
                if 'total_filename_words' not in self.stats:
                    self.stats['total_filename_words'] = len(self.filename_index)
                    
                logger.info(f"Index stats: {self.stats}")
            
            logger.info(f"Index successfully loaded from {filepath}")
            return True
            
        except pickle.UnpicklingError as pe:
            logger.error(f"Pickle loading error - corrupted file: {pe}")
            return False
        except MemoryError as me:
            logger.error(f"Memory error loading index - file too large: {me}")
            return False
        except Exception as e:
            logger.error(f"Failed to load index: {e}", exc_info=True)
            return False
    
    def add_file_turbo(self, file_path: str, content: str, metadata: Dict[str, Any]) -> bool:
        """Ultra-fast file indexing with filename support and logging"""
        try:
            with self._lock:
                # Extract content keywords (ultra-fast)
                content_words = set()
                content_lower = content.lower()
                
                # Split by common delimiters
                for delimiter in [' ', '\n', '\t', '.', ',', ';', ':', '!', '?', '"', "'", '(', ')', '[', ']', '{', '}']:
                    content_lower = content_lower.replace(delimiter, ' ')
                
                # Extract meaningful words from content (length 2-50)
                for word in content_lower.split():
                    word = word.strip()
                    if 2 <= len(word) <= 50 and word.isalnum():
                        content_words.add(word)
                
                # Extract filename keywords
                filename = Path(file_path).stem.lower()  # Filename without extension
                filename_words = set()
                
                # Split filename by common separators
                for delimiter in [' ', '_', '-', '.', '(', ')', '[', ']']:
                    filename = filename.replace(delimiter, ' ')
                
                # Extract meaningful words from filename
                for word in filename.split():
                    word = word.strip()
                    if 2 <= len(word) <= 50 and word.isalnum():
                        filename_words.add(word)
                
                content_keywords = list(content_words)
                filename_keywords = list(filename_words)
                
                # Store in memory index
                self.index_data[file_path] = {
                    'content': content,
                    'metadata': metadata,
                    'keywords': content_keywords,
                    'filename_keywords': filename_keywords,
                    'indexed_time': datetime.now()
                }
                
                # Update content word index
                for word in content_keywords:
                    if word not in self.word_index:
                        self.word_index[word] = set()
                    self.word_index[word].add(file_path)
                
                # Update filename word index
                for word in filename_keywords:
                    if word not in self.filename_index:
                        self.filename_index[word] = set()
                    self.filename_index[word].add(file_path)
                
                # Update stats
                self.stats['total_files'] = len(self.index_data)
                self.stats['total_words'] = len(self.word_index)
                self.stats['total_filename_words'] = len(self.filename_index)
                self.stats['last_update'] = datetime.now()
                
                # Log successful indexing
                ext = metadata.get('extension', Path(file_path).suffix.lower())
                logger.debug(f"Indexed {ext} file: {Path(file_path).name} - Content: {len(content_keywords)} words, Filename: {len(filename_keywords)} words")
                
                return True
                
        except Exception as e:
            logger.error(f"Turbo indexing failed for {file_path}: {e}")
            return False
    
    def search_turbo(self, query: str, search_mode: str = 'both', max_results: int = 100) -> List[Dict[str, Any]]:
        """Ultra-fast search with filename and content support
        
        Args:
            query: Search query
            search_mode: 'content', 'filename', or 'both'
            max_results: Maximum number of results
        """
        if not query or not self.index_data:
            return []
        
        try:
            with self._lock:
                query_words = [word.lower().strip() for word in query.split() if len(word.strip()) >= 2]
                if not query_words:
                    return []
                
                matching_files = set()
                
                # Search in content if requested
                if search_mode in ['content', 'both']:
                    content_matches = self._search_in_index(query_words, self.word_index)
                    matching_files.update(content_matches)
                
                # Search in filenames if requested
                if search_mode in ['filename', 'both']:
                    filename_matches = self._search_in_index(query_words, self.filename_index)
                    matching_files.update(filename_matches)
                
                if not matching_files:
                    return []
                
                # Build results with match type information
                results = []
                for file_path in list(matching_files)[:max_results]:
                    if file_path in self.index_data:
                        file_data = self.index_data[file_path]
                        
                        # Determine match types
                        content_match = self._file_matches_query(query_words, file_data.get('keywords', []))
                        filename_match = self._file_matches_query(query_words, file_data.get('filename_keywords', []))
                        
                        # Calculate relevance score
                        relevance = 0
                        if content_match:
                            relevance += len([w for w in query_words if w in file_data.get('keywords', [])])
                        if filename_match:
                            relevance += len([w for w in query_words if w in file_data.get('filename_keywords', [])]) * 2  # Filename matches get higher score
                        
                        # Create match type indicator
                        match_type = []
                        if filename_match:
                            match_type.append("Filename")
                        if content_match:
                            match_type.append("Content")
                        match_indicator = " + ".join(match_type)
                        
                        # Find matching snippets
                        content = file_data['content']
                        snippet = self._extract_snippet(content, query_words)
                        
                        results.append({
                            'path': file_path,
                            'filename': Path(file_path).name,
                            'snippet': snippet,
                            'metadata': file_data['metadata'],
                            'relevance': relevance,
                            'match_type': match_indicator,
                            'filename_match': filename_match,
                            'content_match': content_match
                        })
                
                # Sort by relevance (filename matches first, then content matches)
                results.sort(key=lambda x: (x['relevance'], x['filename_match']), reverse=True)
                return results
                
        except Exception as e:
            logger.error(f"Turbo search failed: {e}")
            return []
    
    def _search_in_index(self, query_words: List[str], index: Dict[str, set]) -> set:
        """Search for query words in a specific index"""
        matching_files = None
        for word in query_words:
            if word in index:
                word_files = index[word]
                if matching_files is None:
                    matching_files = word_files.copy()
                else:
                    matching_files &= word_files
            else:
                # If any word is not found, no results
                return set()
        
        return matching_files or set()
    
    def _file_matches_query(self, query_words: List[str], file_keywords: List[str]) -> bool:
        """Check if file keywords match all query words"""
        file_keywords_set = set(file_keywords)
        return all(word in file_keywords_set for word in query_words)
    
    def _rebuild_filename_index(self):
        """Rebuild filename index from existing data"""
        self.filename_index = {}
        for file_path, file_data in self.index_data.items():
            filename_keywords = file_data.get('filename_keywords')
            if not filename_keywords:
                # Extract filename keywords for old data
                filename = Path(file_path).stem.lower()
                filename_words = set()
                
                for delimiter in [' ', '_', '-', '.', '(', ')', '[', ']']:
                    filename = filename.replace(delimiter, ' ')
                
                for word in filename.split():
                    word = word.strip()
                    if 2 <= len(word) <= 50 and word.isalnum():
                        filename_words.add(word)
                
                filename_keywords = list(filename_words)
                file_data['filename_keywords'] = filename_keywords
            
            # Update filename index
            for word in filename_keywords:
                if word not in self.filename_index:
                    self.filename_index[word] = set()
                self.filename_index[word].add(file_path)
    
    def _extract_snippet(self, content: str, query_words: List[str], max_length: int = 200) -> str:
        """Extract relevant snippet from content"""
        content_lower = content.lower()
        
        # Find first occurrence of any query word
        best_pos = -1
        for word in query_words:
            pos = content_lower.find(word)
            if pos != -1 and (best_pos == -1 or pos < best_pos):
                best_pos = pos
        
        if best_pos == -1:
            return content[:max_length] + "..." if len(content) > max_length else content
        
        # Extract snippet around the found position
        start = max(0, best_pos - max_length // 2)
        end = min(len(content), start + max_length)
        
        snippet = content[start:end]
        if start > 0:
            snippet = "..." + snippet
        if end < len(content):
            snippet = snippet + "..."
        
        return snippet
    
    def clear_index(self):
        """Clear all indexed data"""
        with self._lock:
            self.index_data.clear()
            self.word_index.clear()
            self.filename_index.clear()
            self.stats = {'total_files': 0, 'total_words': 0, 'total_filename_words': 0, 'last_update': None}
            
            # Clear database
            try:
                with sqlite3.connect(self.db_path) as conn:
                    conn.execute("DELETE FROM files")
                    conn.commit()
            except Exception as e:
                logger.error(f"Failed to clear database: {e}")
    
    def get_stats(self) -> Dict[str, Any]:
        """Get indexing statistics"""
        with self._lock:
            return self.stats.copy()


class TurboIndexingWorker(QThread):
    """Ultra-fast indexing worker thread with time estimation"""
    
    progress_updated = pyqtSignal(int, str, str)  # progress, message, time_estimate
    indexing_completed = pyqtSignal(int, str)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, folders: List[str], extensions: List[str], 
                 search_engine: TurboSearchEngine, extractor: TurboFileExtractor):
        super().__init__()
        self.folders = folders
        self.extensions = set(ext.lower() for ext in extensions)
        self.search_engine = search_engine
        self.extractor = extractor
        self.should_stop = False
        self.start_time = None
        self.processed_files = 0
        
    def stop(self):
        self.should_stop = True
        
    def run(self):
        """Ultra-fast indexing process with time estimation and image file support"""
        try:
            self.start_time = time.time()
            self.processed_files = 0
            
            # Find all files first
            all_files = []
            logger.info("Starting file discovery...")
            
            for folder in self.folders:
                if self.should_stop:
                    return
                logger.info(f"Scanning folder: {folder}")
                folder_files = self._find_files_fast(folder)
                all_files.extend(folder_files)
                logger.info(f"Found {len(folder_files)} files in {folder}")
            
            if not all_files:
                self.indexing_completed.emit(0, "No files found to index")
                return
            
            total_files = len(all_files)
            indexed_count = 0
            
            # Log file type breakdown
            file_types = {}
            for file_path in all_files:
                ext = Path(file_path).suffix.lower()
                file_types[ext] = file_types.get(ext, 0) + 1
            
            logger.info(f"File type breakdown:")
            for ext, count in sorted(file_types.items()):
                logger.info(f"  {ext}: {count} files")
            
            self.progress_updated.emit(0, f"Found {total_files} files. Starting turbo indexing...", "Calculating...")
            
            # Process files with thread pool
            batch_size = min(50, max(1, total_files // 20))  # Adaptive batch size
            
            with ThreadPoolExecutor(max_workers=4) as executor:
                for i in range(0, total_files, batch_size):
                    if self.should_stop:
                        return
                    
                    batch = all_files[i:i + batch_size]
                    futures = [executor.submit(self._process_file, file_path) for file_path in batch]
                    
                    for future in as_completed(futures):
                        if self.should_stop:
                            return
                        
                        try:
                            if future.result():
                                indexed_count += 1
                                self.processed_files += 1
                        except Exception as e:
                            logger.warning(f"File processing error: {e}")
                        
                        # Enhanced time estimation with more frequent updates
                        elapsed_time = time.time() - self.start_time
                        progress = int((indexed_count / total_files) * 100)
                        
                        # More accurate time estimation
                        if indexed_count > 5 and elapsed_time > 2:  # More conservative threshold
                            files_per_second = indexed_count / elapsed_time
                            remaining_files = total_files - indexed_count
                            estimated_remaining = remaining_files / files_per_second
                            time_estimate = self._format_time(estimated_remaining)
                            
                            # Add processing speed info
                            speed_info = f"{files_per_second:.1f} files/sec"
                        elif indexed_count > 0:
                            time_estimate = "Calculating..."
                            speed_info = "Analyzing speed..."
                        else:
                            time_estimate = "Starting..."
                            speed_info = "Initializing..."
                        
                        # Enhanced progress message with speed
                        progress_msg = f"Processing: {indexed_count}/{total_files} files ({speed_info})"
                        
                        # Update progress with enhanced information
                        self.progress_updated.emit(progress, 
                                                 progress_msg,
                                                 f"Est. remaining: {time_estimate}")
            
            if not self.should_stop:
                total_time = time.time() - self.start_time
                logger.info(f"Indexing completed: {indexed_count}/{total_files} files in {self._format_time(total_time)}")
                self.indexing_completed.emit(indexed_count, 
                                           f"Turbo indexing complete! {indexed_count} files indexed in {self._format_time(total_time)}.")
                
        except Exception as e:
            logger.error(f"Indexing error: {e}")
            self.error_occurred.emit(f"Indexing failed: {str(e)}")
    
    def _format_time(self, seconds: float) -> str:
        """Format time in human readable format with enhanced precision"""
        if seconds < 0:
            return "0s"
        elif seconds < 10:
            return f"{seconds:.1f}s"
        elif seconds < 60:
            return f"{int(seconds)}s"
        elif seconds < 3600:
            minutes = int(seconds // 60)
            secs = int(seconds % 60)
            if secs > 0:
                return f"{minutes}m {secs}s"
            else:
                return f"{minutes}m"
        else:
            hours = int(seconds // 3600)
            minutes = int((seconds % 3600) // 60)
            if minutes > 0:
                return f"{hours}h {minutes}m"
            else:
                return f"{hours}h"
    
    def _find_files_fast(self, folder: str) -> List[str]:
        """Ultra-fast file discovery with detailed logging and image file support"""
        files = []
        file_type_count = {}
        image_files = []
        
        try:
            folder_path = Path(folder)
            if not folder_path.exists():
                logger.warning(f"Folder does not exist: {folder}")
                return files
            
            logger.info(f"Scanning folder: {folder}")
            
            for file_path in folder_path.rglob('*'):
                if self.should_stop:
                    break
                
                if file_path.is_file():
                    ext = file_path.suffix.lower()
                    
                    # Count file types
                    file_type_count[ext] = file_type_count.get(ext, 0) + 1
                    
                    if (ext in self.extensions and
                        not any(part.startswith('.') for part in file_path.parts[len(folder_path.parts):])):
                        files.append(str(file_path))
                        
                        # Track image files specifically
                        if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg']:
                            image_files.append(str(file_path))
                        
                        # Limit for ultra speed
                        if len(files) >= 10000:
                            logger.warning(f"Reached 10,000 file limit in {folder}")
                            break
            
            # Log file type statistics
            logger.info(f"Found {len(files)} supported files in {folder}")
            if image_files:
                logger.info(f"Found {len(image_files)} image files")
            
            for ext, count in sorted(file_type_count.items()):
                if ext in self.extensions:
                    icon = "Image" if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg'] else "File"
                    logger.info(f"  {icon} {ext}: {count} files")
                        
        except Exception as e:
            logger.warning(f"Error scanning {folder}: {e}")
        
        return files
    
    def _process_file(self, file_path: str) -> bool:
        """Process single file for indexing with detailed logging"""
        try:
            result = self.extractor.extract_fast(file_path)
            if result['success']:
                # Log successful processing
                ext = Path(file_path).suffix.lower()
                content_preview = result['content'][:100].replace('\n', ' ')
                logger.info(f"Successfully processed {ext} file: {Path(file_path).name} - Content: {content_preview}...")
                
                return self.search_engine.add_file_turbo(
                    file_path, result['content'], result['metadata']
                )
            else:
                # For unsupported files, still try to index by filename
                logger.warning(f"Content extraction failed, but indexing filename: {file_path}")
                try:
                    path = Path(file_path)
                    stat = path.stat()
                    fallback_content = f"File: {path.name}"
                    fallback_metadata = {
                        'filename': path.name,
                        'extension': path.suffix.lower(),
                        'size': stat.st_size,
                        'modified': datetime.fromtimestamp(stat.st_mtime),
                        'created': datetime.fromtimestamp(stat.st_ctime)
                    }
                    return self.search_engine.add_file_turbo(
                        file_path, fallback_content, fallback_metadata
                    )
                except Exception as e:
                    logger.error(f"Failed to create fallback index for {file_path}: {e}")
                    return False
        except Exception as e:
            logger.warning(f"Failed to process {file_path}: {e}")
        return False


class IndexLoadWorker(QThread):
    """Background thread for safe index loading with progress tracking"""
    
    # Signals for communication with main thread
    progress_signal = pyqtSignal(str)  # Progress message
    finished_signal = pyqtSignal(bool, str)  # Success/failure, message
    file_size_signal = pyqtSignal(int)  # File size in bytes
    
    def __init__(self, file_path: str, search_engine):
        super().__init__()
        self.file_path = file_path
        self.search_engine = search_engine
        self.should_stop = False
        
    def stop(self):
        """Stop the loading process"""
        self.should_stop = True
        
    def run(self):
        """Run the enhanced index loading process in background with comprehensive error handling"""
        try:
            logger.info(f"=== IndexLoadWorker ENHANCED Version Started ===")
            logger.info(f"Target file: {self.file_path}")
            logger.info(f"Search engine: {type(self.search_engine).__name__}")
            
            if self.should_stop:
                logger.info("Loading cancelled before starting")
                self.finished_signal.emit(False, "Loading cancelled before starting")
                return
            
            # Step 1: Enhanced file existence and size checking
            self.progress_signal.emit("Performing comprehensive file checks...")
            
            try:
                if not os.path.exists(self.file_path):
                    error_msg = f"File does not exist: {self.file_path}"
                    logger.error(error_msg)
                    self.finished_signal.emit(False, error_msg)
                    return
                
                if not os.path.isfile(self.file_path):
                    error_msg = f"Path is not a file: {self.file_path}"
                    logger.error(error_msg)
                    self.finished_signal.emit(False, error_msg)
                    return
                    
                # Check file size with detailed reporting
                file_size = os.path.getsize(self.file_path)
                self.file_size_signal.emit(file_size)
                logger.info(f"Index file size: {file_size:,} bytes ({file_size/(1024*1024):.2f} MB)")
                
                if file_size == 0:
                    error_msg = "Index file is empty (0 bytes)"
                    logger.error(error_msg)
                    self.finished_signal.emit(False, error_msg)
                    return
                    
                # Warn about very large files
                if file_size > 100 * 1024 * 1024:  # 100MB
                    logger.warning(f"Very large index file detected: {file_size:,} bytes")
                    self.progress_signal.emit(f"Loading large file ({file_size/(1024*1024):.1f} MB)...")
                elif file_size > 10 * 1024 * 1024:  # 10MB
                    logger.info(f"Large index file: {file_size:,} bytes")
                    self.progress_signal.emit(f"Loading file ({file_size/(1024*1024):.1f} MB)...")
                    
            except OSError as os_error:
                error_msg = f"OS error accessing file: {str(os_error)}"
                logger.error(error_msg, exc_info=True)
                self.finished_signal.emit(False, error_msg)
                return
            except Exception as size_error:
                error_msg = f"Error checking file: {str(size_error)}"
                logger.error(error_msg, exc_info=True)
                self.finished_signal.emit(False, error_msg)
                return

            if self.should_stop:
                logger.info("Loading cancelled after file checks")
                self.finished_signal.emit(False, "Loading cancelled")
                return
            
            # Step 2: Enhanced index loading with comprehensive error handling
            self.progress_signal.emit("Loading and validating index data...")
            logger.info("Starting search engine load_index method...")
            
            load_success = False
            try:
                # Monitor the loading process
                start_time = time.time()
                
                # Call the search engine's load_index method with timeout monitoring
                load_success = self.search_engine.load_index(self.file_path)
                
                load_duration = time.time() - start_time
                logger.info(f"Search engine load_index completed in {load_duration:.2f} seconds")
                logger.info(f"Search engine load_index returned: {load_success}")
                
                if self.should_stop:
                    logger.info("Loading cancelled during index loading")
                    self.finished_signal.emit(False, "Loading cancelled")
                    return
                    
            except Exception as load_error:
                error_msg = f"Exception during index loading: {str(load_error)}"
                logger.error(error_msg, exc_info=True)
                self.finished_signal.emit(False, error_msg)
                return
            
            # Step 3: Validate loading results and provide detailed feedback
            if load_success:
                try:
                    self.progress_signal.emit("Index loaded successfully! Finalizing...")
                    
                    # Quick validation of loaded data
                    if hasattr(self.search_engine, 'index_data') and self.search_engine.index_data:
                        index_count = len(self.search_engine.index_data)
                        logger.info(f"Successfully loaded {index_count:,} index entries")
                    else:
                        logger.warning("Index loaded but no data found in search engine")
                    
                    # Small delay to show final message
                    import time
                    time.sleep(0.3)
                    
                    success_msg = f"Index successfully loaded from: {os.path.basename(self.file_path)}"
                    logger.info(f"=== IndexLoadWorker SUCCESS: {success_msg} ===")
                    self.finished_signal.emit(True, success_msg)
                    
                except Exception as validation_error:
                    logger.warning(f"Index loaded but validation failed: {validation_error}")
                    # Still report success since the main loading worked
                    success_msg = f"Index loaded from: {os.path.basename(self.file_path)} (with minor validation warnings)"
                    self.finished_signal.emit(True, success_msg)
                    
            else:
                error_msg = "Failed to load index data - file may be corrupted, incompatible, or from a different version"
                logger.error(error_msg)
                logger.error("Possible causes: corrupted .pkl file, version mismatch, insufficient memory")
                self.finished_signal.emit(False, error_msg)
                
        except Exception as e:
            error_msg = f"Critical error in IndexLoadWorker: {str(e)}"
            logger.error(error_msg, exc_info=True)
            logger.error("This is a critical error that should not occur during normal operation")
            self.finished_signal.emit(False, error_msg)


class DeepFileX(QMainWindow):
    """DeepFileX - Advanced File Analysis System"""

    def __init__(self):
        super().__init__()

        # Initialize turbo components with proper paths
        # Create data directory in user's AppData folder
        data_dir = Path.home() / 'AppData' / 'Roaming' / 'DeepFileX'
        data_dir.mkdir(parents=True, exist_ok=True)
        db_path = data_dir / 'deepfilex.db'

        self.search_engine = TurboSearchEngine(str(db_path))
        self.extractor = TurboFileExtractor()
        self.indexing_worker = None

        # Settings
        self.settings = QSettings('DeepFileX', 'App')
        self.search_paths = []
        self.dark_mode = False
        
        # UI setup
        self.init_ui()
        
        # 🆕 SmartLinks 수익화 시스템 초기화
        if SMARTLINKS_AVAILABLE:
            self.init_smartlinks()
            logger.info("💰 SmartLinks 수익화 시스템 활성화")
        
        # 🆕 자동 업데이트 시스템 초기화
        if UPDATE_SYSTEM_AVAILABLE:
            self.init_update_system()
            logger.info("🔄 자동 업데이트 시스템 활성화")
        
        self.load_settings()
        self.update_stats()
        
        # Auto-search timer
        self.search_timer = QTimer()
        self.search_timer.setSingleShot(True)
        self.search_timer.timeout.connect(self.perform_search)
    
    def init_ui(self):
        """Initialize enhanced UI"""
        self.setWindowTitle("DeepFileX - Advanced File Analysis System")
        self.setGeometry(100, 100, 1500, 1000)
        
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)
        
        # Title with DeepFileX branding
        title = QLabel("DEEPFILEX - ADVANCED FILE ANALYSIS")
        title.setFont(QFont("Arial", 18, QFont.Weight.Bold))
        title.setAlignment(Qt.AlignmentFlag.AlignCenter)
        title.setStyleSheet("""
            QLabel {
                color: #2c5aa0; 
                padding: 15px; 
                background: linear-gradient(45deg, #e8f4fd, #d1e7f4);
                border-radius: 10px;
                margin: 5px;
            }
        """)
        layout.addWidget(title)
        
        # Enhanced Progress bar with better styling
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setMinimumHeight(25)  # Make it more visible
        self.progress_bar.setTextVisible(True)  # Ensure text is visible
        self.progress_bar.setAlignment(Qt.AlignmentFlag.AlignCenter)  # Center the text
        layout.addWidget(self.progress_bar)
        
        # Control panel
        controls = self.create_turbo_controls()
        layout.addWidget(controls)
        
        # Main area
        main_splitter = QSplitter(Qt.Orientation.Horizontal)
        
        # Results panel
        results_panel = self.create_results_panel()
        main_splitter.addWidget(results_panel)
        
        # Info panel
        info_panel = self.create_info_panel()
        main_splitter.addWidget(info_panel)
        
        main_splitter.setSizes([700, 800])
        layout.addWidget(main_splitter)
        
        # Status bar
        self.status_bar = self.statusBar()
        self.status_bar.showMessage("DeepFileX ready - Advanced File Analysis System")
        
        # Apply turbo styling
        self.apply_turbo_styles()
    
    def create_turbo_controls(self):
        """Create enhanced control panel"""
        group = QGroupBox("DeepFileX Analysis Panel")
        layout = QVBoxLayout(group)
        
        # Folder controls
        folder_layout = QHBoxLayout()
        folder_label = QLabel("Scan Folders:")
        folder_label.setMinimumWidth(120)
        folder_layout.addWidget(folder_label)
        
        self.folder_display = QLineEdit()
        self.folder_display.setReadOnly(True)
        self.folder_display.setPlaceholderText("Add folders for deep file analysis...")
        folder_layout.addWidget(self.folder_display)
        
        add_folder_btn = QPushButton("Add Folder")
        add_folder_btn.setText("Add Folder")
        add_folder_btn.clicked.connect(self.add_folder)
        folder_layout.addWidget(add_folder_btn)
        
        clear_btn = QPushButton("Clear")
        clear_btn.setText("Clear")
        clear_btn.clicked.connect(self.clear_folders)
        folder_layout.addWidget(clear_btn)
        
        layout.addLayout(folder_layout)
        
        # Search controls
        search_layout = QHBoxLayout()
        search_label = QLabel("Diagnose:")
        search_label.setMinimumWidth(120)
        search_layout.addWidget(search_label)
        
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("Enter keywords for file diagnosis...")
        self.search_input.textChanged.connect(self.on_search_changed)
        self.search_input.returnPressed.connect(self.perform_search)
        search_layout.addWidget(self.search_input)
        
        search_btn = QPushButton("Start Diagnosis")
        search_btn.setText("Start Diagnosis")
        search_btn.clicked.connect(self.perform_search)
        search_layout.addWidget(search_btn)
        
        clear_search_btn = QPushButton("Clear")
        clear_search_btn.setText("Clear")
        clear_search_btn.clicked.connect(self.clear_search)
        search_layout.addWidget(clear_search_btn)
        
        layout.addLayout(search_layout)
        
        # Search mode controls
        mode_layout = QHBoxLayout()
        mode_label = QLabel("Scan Mode:")
        mode_label.setMinimumWidth(120)
        mode_layout.addWidget(mode_label)
        
        self.search_mode_combo = QComboBox()
        self.search_mode_combo.addItems([
            "Both (Content + Filename)",
            "Content Only", 
            "Filename Only"
        ])
        self.search_mode_combo.setCurrentIndex(0)  # Default to 'both'
        mode_layout.addWidget(self.search_mode_combo)
        
        mode_layout.addStretch()
        layout.addLayout(mode_layout)
        
        # Action controls
        action_layout = QHBoxLayout()
        
        # File type filter
        type_label = QLabel("Type:")
        type_label.setMinimumWidth(120)
        action_layout.addWidget(type_label)
        
        self.file_type_combo = QComboBox()
        self.file_type_combo.addItems([
            'All Types', '.txt', '.py', '.js', '.html', '.css', 
            '.json', '.xml', '.md', '.log', '.csv', '.pdf',
            '.java', '.c', '.cpp', '.h', '.cs', '.php',
            # Office documents
            '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
            # Archives
            '.zip', '.rar', '.7z', '.tar', '.gz',
            # Korean documents
            '.hwp', '.hwpx',
            # Email
            '.pst', '.eml', '.msg',
            # Images
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.svg'
        ])
        action_layout.addWidget(self.file_type_combo)
        
        action_layout.addStretch()
        
        # Index management buttons
        self.index_btn = QPushButton("Start Scan")
        self.index_btn.setText("Start Scan")
        self.index_btn.clicked.connect(self.start_turbo_indexing)
        # Style will be applied by theme
        action_layout.addWidget(self.index_btn)
        
        # Save index button
        save_btn = QPushButton("Save Index")
        save_btn.setText("Save Index")
        save_btn.clicked.connect(self.save_index)
        action_layout.addWidget(save_btn)
        
        # Load index button
        load_btn = QPushButton("Load Index")
        load_btn.setText("Load Index")
        load_btn.clicked.connect(self.load_index)
        action_layout.addWidget(load_btn)
        
        # Clear index button
        clear_index_btn = QPushButton("Clear Records")
        clear_index_btn.setText("Clear Records")
        clear_index_btn.clicked.connect(self.clear_index)
        action_layout.addWidget(clear_index_btn)
        
        # Dark mode toggle
        self.dark_mode_btn = QPushButton("Dark Mode")
        self.dark_mode_btn.setText("Dark Mode")
        self.dark_mode_btn.clicked.connect(self.toggle_dark_mode)
        action_layout.addWidget(self.dark_mode_btn)
        
        # 🆕 Update check button
        if UPDATE_SYSTEM_AVAILABLE:
            self.update_btn = QPushButton("🔄 업데이트")
            self.update_btn.setText("🔄 업데이트")
            self.update_btn.clicked.connect(self.check_for_updates_manual)
            self.update_btn.setStyleSheet("""
                QPushButton {
                    background-color: #27ae60;
                    color: white;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #229954;
                }
            """)
            action_layout.addWidget(self.update_btn)
        
        layout.addLayout(action_layout)
        
        return group
    
    def create_results_panel(self):
        """Create enhanced results panel"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        self.results_header = QLabel("Diagnosis Results (0)")
        self.results_header.setFont(QFont("Arial", 14, QFont.Weight.Bold))
        layout.addWidget(self.results_header)
        
        self.results_list = QListWidget()
        self.results_list.itemClicked.connect(self.on_result_selected)
        self.results_list.itemDoubleClicked.connect(self.open_file_location)
        layout.addWidget(self.results_list)
        
        return widget
    
    def create_info_panel(self):
        """Create enhanced info panel"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        # Stats group
        stats_group = QGroupBox("Analysis Statistics")
        stats_layout = QVBoxLayout(stats_group)
        
        self.stats_label = QLabel("No scan data")
        self.stats_label.setFont(QFont("Arial", 10))
        stats_layout.addWidget(self.stats_label)
        
        layout.addWidget(stats_group)
        
        # Preview group
        preview_group = QGroupBox("File Analysis")
        preview_layout = QVBoxLayout(preview_group)
        
        self.preview_label = QLabel("No file selected")
        self.preview_label.setFont(QFont("Arial", 10, QFont.Weight.Bold))
        preview_layout.addWidget(self.preview_label)
        
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        self.preview_text.setFont(QFont("Consolas", 9))
        preview_layout.addWidget(self.preview_text)
        
        layout.addWidget(preview_group)
        
        return widget
    
    def apply_turbo_styles(self):
        """Apply enhanced turbo styling with theme support"""
        if self.dark_mode:
            self.apply_dark_theme()
        else:
            self.apply_light_theme()
    
    def apply_light_theme(self):
        """Apply light theme styling"""
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f8f9fa;
            }
            QGroupBox {
                font-weight: bold;
                font-size: 12px;
                border: 2px solid #dee2e6;
                border-radius: 8px;
                margin-top: 1ex;
                padding-top: 10px;
                background-color: white;
                color: #495057;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                color: #495057;
            }
            QPushButton {
                background-color: #6c757d;
                color: white;
                border: none;
                padding: 8px 16px;
                border-radius: 5px;
                font-weight: bold;
                min-width: 80px;
            }
            QPushButton:hover {
                background-color: #5a6268;
            }
            QPushButton:pressed {
                background-color: #545b62;
            }
            QLineEdit {
                border: 2px solid #dee2e6;
                border-radius: 5px;
                padding: 8px;
                font-size: 12px;
                background-color: white;
                color: #495057;
            }
            QLineEdit:focus {
                border-color: #e74c3c;
            }
            QComboBox {
                border: 2px solid #dee2e6;
                border-radius: 5px;
                padding: 5px;
                background-color: white;
                min-width: 120px;
                color: #495057;
            }
            QComboBox:focus {
                border-color: #e74c3c;
            }
            QLabel {
                color: #495057;
                font-size: 12px;
            }
            QListWidget {
                border: 1px solid #ddd;
                border-radius: 5px;
                background-color: white;
                color: #495057;
                selection-background-color: #e74c3c;
                selection-color: white;
            }
            QListWidget::item {
                padding: 8px;
                border-bottom: 1px solid #eee;
            }
            QListWidget::item:hover {
                background-color: #f8f9fa;
            }
            QListWidget::item:selected {
                background-color: #e74c3c;
                color: white;
            }
            QTextEdit {
                border: 1px solid #ddd;
                border-radius: 5px;
                background-color: #f8f9fa;
                color: #495057;
            }
            QProgressBar {
                border: 2px solid #e74c3c;
                border-radius: 10px;
                text-align: center;
                font-weight: bold;
                font-size: 14px;
                color: #2c3e50;
                background-color: #f8f9fa;
                min-height: 28px;
            }
            QProgressBar::chunk {
                background: linear-gradient(45deg, #e74c3c, #c0392b);
                border-radius: 8px;
            }
        """)
        
        # Apply special styling for Start Scan button in Light Mode
        if hasattr(self, 'index_btn'):
            self.index_btn.setStyleSheet("""
                QPushButton {
                    background: linear-gradient(45deg, #2c5aa0, #1e3d72);
                    color: #000000;
                    font-weight: bold;
                    padding: 10px 20px;
                    border: 2px solid #2c5aa0;
                    border-radius: 8px;
                    font-size: 14px;
                    min-width: 120px;
                }
                QPushButton:hover {
                    background: linear-gradient(45deg, #1e3d72, #143058);
                    color: #000000;
                    border: 2px solid #3366cc;
                }
                QPushButton:pressed {
                    border: 2px solid #1a2a4a;
                    color: #000000;
                }
            """)
    
    def apply_dark_theme(self):
        """Apply dark theme styling"""
        self.setStyleSheet("""
            QMainWindow {
                background-color: #2b2b2b;
                color: #ffffff;
            }
            QGroupBox {
                font-weight: bold;
                font-size: 12px;
                border: 2px solid #555555;
                border-radius: 8px;
                margin-top: 1ex;
                padding-top: 10px;
                background-color: #3c3c3c;
                color: #ffffff;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                color: #ffffff;
            }
            QPushButton {
                background-color: #555555;
                color: white;
                border: none;
                padding: 8px 16px;
                border-radius: 5px;
                font-weight: bold;
                min-width: 80px;
            }
            QPushButton:hover {
                background-color: #666666;
            }
            QPushButton:pressed {
                background-color: #444444;
            }
            QLineEdit {
                border: 2px solid #555555;
                border-radius: 5px;
                padding: 8px;
                font-size: 12px;
                background-color: #3c3c3c;
                color: #ffffff;
            }
            QLineEdit:focus {
                border-color: #e74c3c;
            }
            QComboBox {
                border: 2px solid #555555;
                border-radius: 5px;
                padding: 5px;
                background-color: #3c3c3c;
                min-width: 120px;
                color: #ffffff;
            }
            QComboBox:focus {
                border-color: #e74c3c;
            }
            QComboBox::drop-down {
                background-color: #555555;
            }
            QComboBox QAbstractItemView {
                background-color: #3c3c3c;
                color: #ffffff;
                selection-background-color: #e74c3c;
                selection-color: #ffffff;
            }
            QLabel {
                color: #ffffff;
                font-size: 12px;
            }
            QListWidget {
                border: 1px solid #555555;
                border-radius: 5px;
                background-color: #3c3c3c;
                color: #ffffff;
                selection-background-color: #e74c3c;
                selection-color: #ffffff;
            }
            QListWidget::item {
                padding: 8px;
                border-bottom: 1px solid #555555;
                color: #ffffff;
            }
            QListWidget::item:hover {
                background-color: #4a4a4a;
                color: #ffffff;
            }
            QListWidget::item:selected {
                background-color: #e74c3c;
                color: #ffffff;
            }
            QTextEdit {
                border: 1px solid #555555;
                border-radius: 5px;
                background-color: #3c3c3c;
                color: #ffffff;
            }
            QProgressBar {
                border: 2px solid #e74c3c;
                border-radius: 10px;
                text-align: center;
                font-weight: bold;
                font-size: 14px;
                color: #ffffff;
                background-color: #3c3c3c;
                min-height: 28px;
            }
            QProgressBar::chunk {
                background: linear-gradient(45deg, #e74c3c, #c0392b);
                border-radius: 8px;
            }
        """)
        
        # Apply special styling for Start Scan button in Dark Mode
        if hasattr(self, 'index_btn'):
            self.index_btn.setStyleSheet("""
                QPushButton {
                    background: linear-gradient(45deg, #2c5aa0, #1e3d72);
                    color: #ffffff;
                    font-weight: bold;
                    padding: 10px 20px;
                    border: 2px solid #2c5aa0;
                    border-radius: 8px;
                    font-size: 14px;
                    min-width: 120px;
                }
                QPushButton:hover {
                    background: linear-gradient(45deg, #1e3d72, #143058);
                    color: #ffffff;
                    border: 2px solid #3366cc;
                }
                QPushButton:pressed {
                    border: 2px solid #1a2a4a;
                    color: #ffffff;
                }
            """)
    
    def toggle_dark_mode(self):
        """Toggle between dark and light mode"""
        self.dark_mode = not self.dark_mode
        self.dark_mode_btn.setText("Light Mode" if self.dark_mode else "Dark Mode")
        self.apply_turbo_styles()
        self.save_settings()
        
        # Force refresh of all widgets to apply new theme
        self.update()
        if hasattr(self, 'results_list'):
            self.results_list.update()
        if hasattr(self, 'preview_text'):
            self.preview_text.update()
        
        # Log theme change
        theme_name = "Dark" if self.dark_mode else "Light"
        logger.info(f"Theme changed to {theme_name} mode")
        self.status_bar.showMessage(f"Switched to {theme_name} mode")
    
    def add_folder(self):
        """Add folder for searching"""
        folder = QFileDialog.getExistingDirectory(self, "Select Folder to Search")
        if folder and folder not in self.search_paths:
            self.search_paths.append(folder)
            self.update_folder_display()
            self.save_settings()
    
    def clear_folders(self):
        """Clear all search folders"""
        self.search_paths.clear()
        self.update_folder_display()
        self.save_settings()
    
    def update_folder_display(self):
        """Update folder display"""
        if self.search_paths:
            display_text = "; ".join(self.search_paths)
            if len(display_text) > 100:
                display_text = display_text[:97] + "..."
            self.folder_display.setText(display_text)
        else:
            self.folder_display.setText("")
    
    def save_index(self):
        """Save current index to file with dedicated indexes folder"""
        if not self.search_engine.index_data:
            QMessageBox.warning(self, "No Data", "No scan data to save.")
            return
        
        # Create indexes folder in user's AppData folder
        data_dir = Path.home() / 'AppData' / 'Roaming' / 'DeepFileX'
        indexes_folder = data_dir / "indexes"
        os.makedirs(indexes_folder, exist_ok=True)

        # Generate folder names for filename
        folder_names = self._generate_folder_names_for_filename()

        # Default filename with folder names and timestamp
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        if folder_names:
            default_filename = f"deepfilex_index_{folder_names}_{timestamp}.pkl"
        else:
            default_filename = f"deepfilex_index_{timestamp}.pkl"
        
        default_path = os.path.join(indexes_folder, default_filename)
        
        file_path, _ = QFileDialog.getSaveFileName(
            self, "Save Index", 
            default_path,
            "Pickle files (*.pkl);;All files (*.*)"
        )
        
        if file_path:
            if self.search_engine.save_index(file_path):
                QMessageBox.information(self, "Index Saved", 
                                      f"Index successfully saved to:\n{file_path}")
            else:
                QMessageBox.critical(self, "Save Failed", "Failed to save index.")
    
    def _generate_folder_names_for_filename(self):
        """Generate a safe folder name string for use in filenames"""
        if not self.search_paths:
            return ""
        
        folder_names = []
        for path in self.search_paths:
            try:
                # Get the last directory name from the path
                folder_name = os.path.basename(os.path.normpath(path))
                
                # Clean the name for use in filename (remove invalid characters)
                safe_name = self._make_filename_safe(folder_name)
                
                if safe_name and safe_name not in folder_names:
                    folder_names.append(safe_name)
                    
            except Exception as e:
                logger.warning(f"Error processing folder name for {path}: {e}")
                continue
        
        if not folder_names:
            return ""
        
        # Join folder names with hyphens
        combined_name = "-".join(folder_names)
        
        # Limit total length to avoid extremely long filenames
        max_length = 50
        if len(combined_name) > max_length:
            # Truncate and add indicator
            combined_name = combined_name[:max_length-3] + "..."
        
        return combined_name
    
    def _make_filename_safe(self, name):
        """Make a string safe for use in filenames"""
        if not name:
            return ""
        
        # Remove or replace invalid filename characters
        invalid_chars = '<>:"/\\|?*'
        safe_name = name
        
        for char in invalid_chars:
            safe_name = safe_name.replace(char, "")
        
        # Replace spaces and dots with underscores
        safe_name = safe_name.replace(" ", "_").replace(".", "_")
        
        # Remove multiple consecutive underscores
        while "__" in safe_name:
            safe_name = safe_name.replace("__", "_")
        
        # Remove leading/trailing underscores
        safe_name = safe_name.strip("_")
        
        # Limit individual folder name length
        if len(safe_name) > 20:
            safe_name = safe_name[:20]
        
        return safe_name
    
    def load_index(self):
        """Load index from file with enhanced safety and PyInstaller compatibility"""
        try:
            logger.info("=== ENHANCED Load Index Process Started ===")
            
            # Step 1: Determine indexes folder path with maximum robustness
            indexes_folder = None
            
            try:
                # Primary: User AppData folder
                user_home = os.path.expanduser('~')
                data_dir = os.path.join(user_home, 'AppData', 'Roaming', 'DeepFileX')
                indexes_folder = os.path.join(data_dir, 'indexes')
                
                # Create folder if it doesn't exist
                os.makedirs(indexes_folder, exist_ok=True)
                logger.info(f"Using primary indexes folder: {indexes_folder}")
                
            except Exception as primary_error:
                logger.warning(f"Could not use primary indexes folder: {primary_error}")
                
                try:
                    # Fallback: Current directory + indexes
                    indexes_folder = os.path.join(os.getcwd(), 'indexes')
                    os.makedirs(indexes_folder, exist_ok=True)
                    logger.info(f"Using fallback indexes folder: {indexes_folder}")
                    
                except Exception as fallback_error:
                    logger.warning(f"Could not create fallback folder: {fallback_error}")
                    # Final fallback: current directory
                    indexes_folder = os.getcwd()
                    logger.info(f"Using current directory as final fallback: {indexes_folder}")
            
            # Step 2: Enhanced file dialog with multiple fallback methods
            logger.info(f"Opening file dialog in: {indexes_folder}")
            
            file_path = None
            dialog_success = False
            
            # Method 1: Try with DontUseNativeDialog (PyInstaller compatible)
            try:
                logger.info("Attempting Method 1: Non-native file dialog")
                from PyQt6.QtWidgets import QFileDialog
                
                dialog_options = QFileDialog.Option.DontUseNativeDialog
                
                file_path, file_filter = QFileDialog.getOpenFileName(
                    self,
                    "Load DeepFileX Index File",
                    indexes_folder,
                    "DeepFileX Index Files (*.pkl);;All Files (*.*)",
                    options=dialog_options
                )
                
                if file_path:  # User selected a file
                    logger.info(f"Method 1 SUCCESS: path='{file_path}', filter='{file_filter}'")
                    dialog_success = True
                else:
                    logger.info("Method 1: User cancelled dialog")
                    return  # User cancelled, this is not an error
                
            except Exception as method1_error:
                logger.warning(f"Method 1 failed: {method1_error}")
                
                # Method 2: Try with native dialog (fallback)
                try:
                    logger.info("Attempting Method 2: Native file dialog")
                    
                    file_path, file_filter = QFileDialog.getOpenFileName(
                        self,
                        "Load DeepFileX Index File",
                        indexes_folder,
                        "DeepFileX Index Files (*.pkl);;All Files (*.*)"
                    )
                    
                    if file_path:
                        logger.info(f"Method 2 SUCCESS: path='{file_path}', filter='{file_filter}'")
                        dialog_success = True
                    else:
                        logger.info("Method 2: User cancelled dialog")
                        return
                        
                except Exception as method2_error:
                    logger.error(f"Method 2 failed: {method2_error}")
                    
                    # Method 3: Manual file input as last resort
                    try:
                        logger.info("Attempting Method 3: Manual file input")
                        from PyQt6.QtWidgets import QInputDialog
                        
                        file_path, ok = QInputDialog.getText(
                            self, 
                            'Load Index File', 
                            f'Please enter the full path to your index file:\n(Default folder: {indexes_folder})',
                            text=os.path.join(indexes_folder, 'deepfilex_index_*.pkl')
                        )
                        
                        if ok and file_path.strip():
                            file_path = file_path.strip()
                            logger.info(f"Method 3 SUCCESS: path='{file_path}'")
                            dialog_success = True
                        else:
                            logger.info("Method 3: User cancelled")
                            return
                            
                    except Exception as method3_error:
                        logger.error(f"All file dialog methods failed. Method 3 error: {method3_error}", exc_info=True)
                        QMessageBox.critical(
                            self, 
                            "File Dialog Error", 
                            f"Could not open file selection dialog:\n\n"
                            f"Method 1 Error: {str(method1_error)}\n"
                            f"Method 2 Error: {str(method2_error)}\n"
                            f"Method 3 Error: {str(method3_error)}\n\n"
                            f"Please try restarting the application or check the log file."
                        )
                        return
            
            # Step 3: Enhanced file validation with detailed checking
            if not file_path or not file_path.strip():
                logger.info("No file selected - user cancelled dialog or provided empty path")
                return
                
            file_path = file_path.strip()
            logger.info(f"Selected file path: '{file_path}'")
                
            # Check file existence with detailed error reporting
            if not os.path.exists(file_path):
                error_msg = f"The selected file does not exist:\n{file_path}\n\nPlease check if:\n• The file exists\n• The path is correct\n• You have read permissions"
                logger.error(f"File not found: {file_path}")
                QMessageBox.critical(self, "File Not Found", error_msg)
                return
            
            # Check if it's actually a file (not a directory)
            if not os.path.isfile(file_path):
                error_msg = f"The selected path is not a file:\n{file_path}\n\nPlease select a valid .pkl file."
                logger.error(f"Path is not a file: {file_path}")
                QMessageBox.critical(self, "Invalid File", error_msg)
                return
            
            # Check file permissions
            try:
                with open(file_path, 'rb') as test_file:
                    test_file.read(1)  # Try to read 1 byte
                    logger.info("File access test passed")
            except PermissionError as perm_error:
                error_msg = f"Permission denied accessing file:\n{file_path}\n\nPlease check file permissions."
                logger.error(f"Permission error: {perm_error}")
                QMessageBox.critical(self, "Permission Error", error_msg)
                return
            except Exception as access_error:
                error_msg = f"Cannot access file:\n{file_path}\n\nError: {str(access_error)}"
                logger.error(f"File access error: {access_error}")
                QMessageBox.critical(self, "File Access Error", error_msg)
                return
            
            # Step 4: Enhanced background loading with better error handling
            logger.info(f"Starting enhanced background loading of: {file_path}")
            
            try:
                # Create progress dialog with enhanced options
                from PyQt6.QtWidgets import QProgressDialog
                
                self.load_progress_dialog = QProgressDialog(
                    "Preparing to load index...", "Cancel", 0, 0, self
                )
                self.load_progress_dialog.setWindowTitle("DeepFileX - Loading Index")
                self.load_progress_dialog.setWindowModality(Qt.WindowModality.WindowModal)
                self.load_progress_dialog.setMinimumWidth(450)
                self.load_progress_dialog.setMinimumHeight(120)
                
                # Ensure dialog is visible
                self.load_progress_dialog.show()
                QApplication.processEvents()  # Force UI update
                
                logger.info("Progress dialog created and shown")
                
                # Create and configure the loading worker with enhanced error handling
                self.index_load_worker = IndexLoadWorker(file_path, self.search_engine)
                
                # Connect signals with enhanced error handling
                try:
                    self.index_load_worker.progress_signal.connect(self.load_progress_dialog.setLabelText)
                    self.index_load_worker.finished_signal.connect(self.on_index_load_finished)
                    self.index_load_worker.file_size_signal.connect(self.on_index_file_size)
                    
                    # Connect cancel button
                    self.load_progress_dialog.canceled.connect(self.on_index_load_cancelled)
                    
                    logger.info("All signals connected successfully")
                    
                except Exception as signal_error:
                    logger.error(f"Error connecting signals: {signal_error}", exc_info=True)
                    if hasattr(self, 'load_progress_dialog'):
                        self.load_progress_dialog.close()
                    QMessageBox.critical(
                        self, 
                        "Signal Connection Error", 
                        f"Could not set up background loading:\n\n{str(signal_error)}\n\nPlease try again."
                    )
                    return
                
                # Start the worker with error handling
                try:
                    self.index_load_worker.start()
                    logger.info("Index loading worker started successfully")
                    
                except Exception as start_error:
                    logger.error(f"Error starting worker: {start_error}", exc_info=True)
                    if hasattr(self, 'load_progress_dialog'):
                        self.load_progress_dialog.close()
                    QMessageBox.critical(
                        self, 
                        "Worker Start Error", 
                        f"Could not start background loading:\n\n{str(start_error)}\n\nPlease try again."
                    )
                    return
                    
            except Exception as dialog_setup_error:
                logger.error(f"Error setting up progress dialog: {dialog_setup_error}", exc_info=True)
                QMessageBox.critical(
                    self, 
                    "Dialog Setup Error", 
                    f"Could not create progress dialog:\n\n{str(dialog_setup_error)}\n\nPlease try again."
                )
                return
            
        except Exception as e:
            logger.error(f"CRITICAL ERROR in enhanced load_index: {e}", exc_info=True)
            QMessageBox.critical(
                self, 
                "Critical Error", 
                f"A critical error occurred during the index loading process:\n\n{str(e)}\n\n"
                f"This error has been logged. Please check the log file at:\n{log_file}\n\n"
                f"If this problem persists, please restart the application."
            )
    
    def on_index_file_size(self, file_size: int):
        """Handle file size information from worker"""
        size_mb = file_size / (1024 * 1024)
        if size_mb > 10:
            self.load_progress_dialog.setLabelText(f"Loading large index file ({size_mb:.1f} MB)...")
    
    def on_index_load_cancelled(self):
        """Handle user cancellation of index loading with comprehensive cleanup"""
        try:
            logger.info("=== Index Loading Cancellation Handler Started ===")
            
            # Step 1: Stop the worker thread safely
            if hasattr(self, 'index_load_worker') and self.index_load_worker:
                try:
                    logger.info("Attempting to stop worker thread gracefully")
                    self.index_load_worker.stop()
                    
                    # Wait for graceful stop with timeout
                    stopped_gracefully = self.index_load_worker.wait(3000)  # Wait up to 3 seconds
                    
                    if stopped_gracefully:
                        logger.info("Worker stopped gracefully")
                    else:
                        logger.warning("Worker did not stop gracefully within 3 seconds, terminating forcefully")
                        
                        try:
                            self.index_load_worker.terminate()
                            terminated = self.index_load_worker.wait(2000)  # Wait up to 2 seconds for termination
                            
                            if terminated:
                                logger.info("Worker terminated successfully")
                            else:
                                logger.error("Worker could not be terminated - this may cause issues")
                                
                        except Exception as terminate_error:
                            logger.error(f"Error during worker termination: {terminate_error}")
                    
                    # Clean up worker reference
                    try:
                        self.index_load_worker.deleteLater()
                        self.index_load_worker = None
                        logger.info("Worker reference cleaned up")
                    except Exception as cleanup_error:
                        logger.warning(f"Error cleaning up worker reference: {cleanup_error}")
                        
                except Exception as worker_error:
                    logger.error(f"Error stopping worker: {worker_error}", exc_info=True)
            else:
                logger.info("No worker thread to stop")
            
            # Step 2: Close progress dialog safely
            try:
                if hasattr(self, 'load_progress_dialog') and self.load_progress_dialog:
                    logger.info("Closing progress dialog")
                    self.load_progress_dialog.close()
                    self.load_progress_dialog.deleteLater()
                    self.load_progress_dialog = None
                    logger.info("Progress dialog closed successfully")
                else:
                    logger.info("No progress dialog to close")
            except Exception as dialog_error:
                logger.warning(f"Error closing progress dialog: {dialog_error}")
            
            # Step 3: Update UI to reflect cancellation
            try:
                # Force UI update
                QApplication.processEvents()
                logger.info("UI updated after cancellation")
            except Exception as ui_error:
                logger.warning(f"Error updating UI: {ui_error}")
            
            logger.info("=== Index Loading Cancellation Handler Completed ===")
            
        except Exception as e:
            logger.error(f"CRITICAL ERROR in on_index_load_cancelled: {e}", exc_info=True)
            # Try to at least close the dialog even if other cleanup fails
            try:
                if hasattr(self, 'load_progress_dialog') and self.load_progress_dialog:
                    self.load_progress_dialog.close()
            except:
                pass
    
    def on_index_load_finished(self, success: bool, message: str):
        """Handle completion of index loading with enhanced error handling and user feedback"""
        try:
            logger.info(f"=== Index Loading Completion Handler ===")
            logger.info(f"Success: {success}, Message: '{message}'")
            
            # Step 1: Clean up progress dialog safely
            try:
                if hasattr(self, 'load_progress_dialog') and self.load_progress_dialog:
                    logger.info("Closing progress dialog")
                    self.load_progress_dialog.close()
                    self.load_progress_dialog.deleteLater()
                    self.load_progress_dialog = None
                    logger.info("Progress dialog closed successfully")
            except Exception as dialog_error:
                logger.warning(f"Error closing progress dialog: {dialog_error}")
            
            # Step 2: Clean up worker thread safely
            try:
                if hasattr(self, 'index_load_worker') and self.index_load_worker:
                    logger.info("Cleaning up worker thread")
                    
                    # Ensure worker is stopped
                    if self.index_load_worker.isRunning():
                        logger.info("Worker still running, stopping it")
                        self.index_load_worker.stop()
                        self.index_load_worker.wait(2000)  # Wait up to 2 seconds
                        
                        if self.index_load_worker.isRunning():
                            logger.warning("Worker didn't stop gracefully, terminating")
                            self.index_load_worker.terminate()
                            self.index_load_worker.wait()
                    
                    self.index_load_worker.deleteLater()
                    self.index_load_worker = None
                    logger.info("Worker thread cleaned up successfully")
                    
            except Exception as worker_error:
                logger.warning(f"Error cleaning up worker: {worker_error}")
            
            # Step 3: Handle success case with comprehensive UI updates
            if success:
                try:
                    logger.info("Processing successful index load")
                    
                    # Update UI stats with error handling
                    try:
                        self.update_stats()
                        logger.info("UI stats updated successfully after index load")
                    except Exception as stats_error:
                        logger.warning(f"Could not update stats after successful load: {stats_error}")
                        # Don't fail the entire process for stats update failure
                    
                    # Force UI refresh
                    try:
                        QApplication.processEvents()
                        logger.info("UI refresh completed")
                    except Exception as ui_error:
                        logger.warning(f"UI refresh error: {ui_error}")
                    
                    # Show enhanced success message
                    try:
                        # Get additional info about the loaded index
                        index_info = ""
                        if hasattr(self.search_engine, 'index_data') and self.search_engine.index_data:
                            file_count = len(self.search_engine.index_data)
                            index_info = f"\n\nIndex contains {file_count:,} files."
                        
                        success_message = (
                            f"{message}{index_info}\n\n"
                            f"✓ Index successfully loaded\n"
                            f"✓ Search engine updated\n"  
                            f"✓ Ready for file diagnosis\n\n"
                            f"You can now perform searches using the loaded index."
                        )
                        
                        QMessageBox.information(
                            self,
                            "Index Loaded Successfully",
                            success_message
                        )
                        logger.info("Success message shown to user")
                        
                    except Exception as msg_error:
                        logger.warning(f"Error showing success message: {msg_error}")
                        # Fallback to simple message
                        QMessageBox.information(
                            self,
                            "Index Loaded",
                            f"{message}\n\nIndex loaded successfully."
                        )
                    
                except Exception as success_error:
                    logger.error(f"Error handling successful load: {success_error}", exc_info=True)
                    # Still show a message to the user
                    QMessageBox.information(
                        self,
                        "Index Loaded",
                        f"{message}\n\n(Note: Some post-loading updates may have failed - see log for details)"
                    )
                    
            # Step 4: Handle failure case with detailed error reporting
            else:
                try:
                    logger.error(f"Processing failed index load: {message}")
                    
                    # Analyze the error message to provide better guidance
                    guidance = ""
                    if "does not exist" in message.lower():
                        guidance = "\n\nSuggestions:\n• Check if the file was moved or deleted\n• Verify the file path is correct"
                    elif "permission" in message.lower():
                        guidance = "\n\nSuggestions:\n• Check file permissions\n• Try running as administrator\n• Close other programs that might be using the file"
                    elif "corrupted" in message.lower() or "incompatible" in message.lower():
                        guidance = "\n\nSuggestions:\n• Try a different index file\n• Re-create the index from scratch\n• Check if the file was created with a compatible version"
                    elif "memory" in message.lower():
                        guidance = "\n\nSuggestions:\n• Close other applications to free memory\n• Try a smaller index file\n• Restart the application"
                    
                    error_message = (
                        f"Failed to load index file.\n\n"
                        f"Error Details:\n{message}{guidance}\n\n"
                        f"You can:\n"
                        f"• Try a different index file\n"
                        f"• Create a new index by scanning folders\n"
                        f"• Check the log file for more details"
                    )
                    
                    QMessageBox.critical(
                        self,
                        "Index Load Failed",
                        error_message
                    )
                    logger.info("Error message shown to user")
                    
                except Exception as error_handling_error:
                    logger.error(f"Error handling failed load: {error_handling_error}", exc_info=True)
                    # Fallback error message
                    QMessageBox.critical(
                        self,
                        "Load Failed",
                        f"Failed to load index:\n\n{message}"
                    )
            
            logger.info("=== Index Loading Completion Handler Finished ===")
            
        except Exception as e:
            logger.error(f"CRITICAL ERROR in on_index_load_finished: {e}", exc_info=True)
            try:
                QMessageBox.critical(
                    self,
                    "Critical Error",
                    f"A critical error occurred after index loading:\n\n{str(e)}\n\n"
                    f"Please restart the application and try again."
                )
            except:
                logger.error("Could not even show error message dialog - system may be unstable")
                pass
    
    def start_turbo_indexing(self):
        """Start or stop turbo indexing with confirmation dialog"""
        if self.indexing_worker and self.indexing_worker.isRunning():
            self.indexing_worker.stop()
            self.indexing_worker.wait()
            self.index_btn.setText("Start Scan")
            self.progress_bar.setVisible(False)
            self.status_bar.showMessage("Indexing stopped")
            return
        
        if not self.search_paths:
            QMessageBox.warning(self, "No Folders", 
                              "Please add folders before starting file scan.")
            return
        
        # Show confirmation dialog with detailed information
        if not self.show_indexing_confirmation():
            return
        
        # Start turbo indexing with all supported file types
        extensions = [
            # Text and code files
            '.txt', '.py', '.js', '.html', '.css', '.json', '.xml', 
            '.md', '.log', '.csv', '.java', '.c', '.cpp', '.h', '.cs',
            '.pdf', '.php', '.rb', '.go', '.rs', '.swift', '.kt',
            '.yaml', '.yml', '.ini', '.cfg', '.conf', '.toml',
            # Office documents
            '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls',
            # Archives (for filename search)
            '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2',
            # Korean documents
            '.hwp', '.hwpx',
            # Email and other formats
            '.pst', '.eml', '.msg', '.rtf',
            # Image files (for filename search)
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg'
        ]
        
        self.indexing_worker = TurboIndexingWorker(
            self.search_paths, extensions, self.search_engine, self.extractor
        )
        
        self.indexing_worker.progress_updated.connect(self.on_progress_update)
        self.indexing_worker.indexing_completed.connect(self.on_indexing_complete)
        self.indexing_worker.error_occurred.connect(self.on_indexing_error)
        
        self.index_btn.setText("Stop")
        
        # Enhanced progress bar setup - ALWAYS visible during indexing
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
        self.progress_bar.setFormat("0% • Preparing to start...")
        
        # Force progress bar styling that won't be overridden
        progress_style = """
            QProgressBar {
                border: 2px solid #e74c3c;
                border-radius: 10px;
                text-align: center;
                font-weight: bold;
                font-size: 14px;
                color: #ffffff;
                background-color: #2c3e50;
                min-height: 30px;
                max-height: 30px;
            }
            QProgressBar::chunk {
                background: linear-gradient(45deg, #e74c3c, #c0392b);
                border-radius: 8px;
            }
        """
        self.progress_bar.setStyleSheet(progress_style)
        
        # Ensure progress bar is updated
        self.progress_bar.update()
        self.progress_bar.repaint()
        
        self.status_bar.showMessage("Initializing insight indexing... Please wait...")
        
        logger.info("Starting insight indexing with enhanced progress tracking")
        self.indexing_worker.start()
    
    def show_indexing_confirmation(self) -> bool:
        """Show indexing confirmation dialog with detailed information"""
        try:
            # Build simple folder summary without deep scanning (no delay)
            folder_summary = "\n".join([
                f"Folder: {Path(path).name} ({path})"
                for path in self.search_paths
            ])
            
            # Create confirmation message without file counting to avoid delay
            confirmation_msg = f"""
TURBO INDEXING CONFIRMATION

FOLDERS TO SCAN:
{folder_summary}

IMPORTANT INFORMATION:
• Processing time depends on computer specifications and data volume
• Large files or slow storage may increase indexing time
• You can stop indexing at any time using the Stop button
• The process will create searchable indexes for faster future searches

SUPPORTED FILE TYPES:
• Text files: .txt, .md, .log, .csv, .json, .xml, .html, .css
• Code files: .py, .js, .java, .c, .cpp, .h, .cs, .php, .rb, .go
• Documents: .pdf, .docx, .doc, .pptx, .ppt, .xlsx, .xls
• Images: .jpg, .jpeg, .png, .gif, .bmp, .tiff, .webp, .ico, .svg
• Archives: .zip, .rar, .7z, .tar, .gz, .bz2
• Other formats: .hwp, .hwpx, .pst, .eml, .msg, .rtf

RECOMMENDATIONS:
• Close other heavy applications for faster processing
• Ensure sufficient disk space for index files
• Let the process complete for best search performance

WHAT HAPPENS NEXT:
• Files will be scanned and content extracted
• Search indexes will be built in memory
• Progress will be shown in real-time
• You'll get a completion notification when done

Do you want to start turbo indexing now?
            """.strip()
            
            # Show confirmation dialog
            msg_box = QMessageBox(self)
            msg_box.setWindowTitle("Turbo Indexing Confirmation")
            msg_box.setText(confirmation_msg)
            msg_box.setIcon(QMessageBox.Icon.Question)
            msg_box.setStandardButtons(QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            msg_box.setDefaultButton(QMessageBox.StandardButton.Yes)
            
            # Customize button text
            yes_button = msg_box.button(QMessageBox.StandardButton.Yes)
            no_button = msg_box.button(QMessageBox.StandardButton.No)
            yes_button.setText("Start Indexing")
            no_button.setText("Cancel")
            
            # Apply theme to message box for better visibility
            if self.dark_mode:
                msg_box.setStyleSheet("""
                    QMessageBox {
                        background-color: #3c3c3c;
                        color: #ffffff;
                        font-size: 13px;
                        font-weight: normal;
                    }
                    QMessageBox QLabel {
                        color: #ffffff;
                        font-size: 13px;
                        background-color: transparent;
                    }
                    QMessageBox QPushButton {
                        background-color: #555555;
                        color: #ffffff;
                        border: none;
                        padding: 8px 16px;
                        border-radius: 5px;
                        font-weight: bold;
                        font-size: 12px;
                        min-width: 80px;
                    }
                    QMessageBox QPushButton:hover {
                        background-color: #666666;
                        color: #ffffff;
                    }
                """)
            else:
                # Light mode styling for better contrast
                msg_box.setStyleSheet("""
                    QMessageBox {
                        background-color: #ffffff;
                        color: #000000;
                        font-size: 13px;
                        font-weight: normal;
                        border: 1px solid #cccccc;
                    }
                    QMessageBox QLabel {
                        color: #000000;
                        font-size: 13px;
                        background-color: transparent;
                        font-weight: normal;
                    }
                    QMessageBox QPushButton {
                        background-color: #e74c3c;
                        color: #ffffff;
                        border: none;
                        padding: 8px 16px;
                        border-radius: 5px;
                        font-weight: bold;
                        font-size: 12px;
                        min-width: 80px;
                    }
                    QMessageBox QPushButton:hover {
                        background-color: #c0392b;
                        color: #ffffff;
                    }
                """)
            
            result = msg_box.exec()
            return result == QMessageBox.StandardButton.Yes
            
        except Exception as e:
            logger.error(f"Error in indexing confirmation: {e}")
            # Fallback to simple confirmation
            reply = QMessageBox.question(
                self, "Start Indexing?", 
                "Start turbo indexing? This may take some time depending on the number of files.",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            return reply == QMessageBox.StandardButton.Yes
    
    def on_progress_update(self, progress, message, time_estimate):
        """Handle progress updates with enhanced time estimation display"""
        # Ensure progress bar is visible and properly updated
        if not self.progress_bar.isVisible():
            self.progress_bar.setVisible(True)
        
        self.progress_bar.setValue(progress)
        
        # Enhanced progress bar format with clearer time display
        if isinstance(time_estimate, str) and time_estimate and time_estimate != "Calculating...":
            if "Est. remaining:" in time_estimate:
                time_part = time_estimate.replace("Est. remaining: ", "")
                format_text = f"{progress}% • {time_part} remaining"
            else:
                format_text = f"{progress}% • {time_estimate}"
        else:
            format_text = f"{progress}% • Calculating time..."
        
        self.progress_bar.setFormat(format_text)
        
        # Enhanced status message
        if progress > 0:
            status_message = f"{message} • {progress}% complete"
        else:
            status_message = message
            
        self.status_bar.showMessage(status_message)
        
        # Force immediate UI update for smoother display
        self.progress_bar.update()
        self.progress_bar.repaint()
        self.status_bar.repaint()
        
        # Log progress for debugging
        logger.debug(f"Progress update: {progress}% - {format_text}")
    
    def on_indexing_complete(self, count, message):
        """Handle indexing completion with detailed statistics"""
        self.index_btn.setText("Start Scan")
        
        # Enhanced progress bar completion
        self.progress_bar.setValue(100)
        self.progress_bar.setFormat("100% • Indexing Complete!")
        
        # Keep progress bar visible for a moment to show completion
        QTimer.singleShot(2000, lambda: self.progress_bar.setVisible(False))
        
        self.status_bar.showMessage(f"{message}")
        self.update_stats()
        
        try:
            # Show detailed completion message with file type breakdown
            stats = self.search_engine.get_stats()
            
            # Safely get statistics with defaults
            total_files = stats.get('total_files', count)
            total_words = stats.get('total_words', 0)
            filename_words = stats.get('total_filename_words', 0)
            
            # Check how many image files were indexed
            indexed_files = self.search_engine.index_data
            image_count = 0
            pdf_count = 0
            office_count = 0
            
            for file_path in indexed_files:
                ext = Path(file_path).suffix.lower()
                if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg']:
                    image_count += 1
                elif ext == '.pdf':
                    pdf_count += 1
                elif ext in ['.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls']:
                    office_count += 1
            
            completion_msg = (f"Scan Complete!\n\n"
                             f"SCAN RESULTS:\n"
                             f"• Total Files Indexed: {total_files:,}\n"
                             f"• Content Words Indexed: {total_words:,}\n"
                             f"• Filename Words Indexed: {filename_words:,}\n\n"
                             f"FILE TYPE BREAKDOWN:\n"
                             f"• Image Files: {image_count:,}\n"
                             f"• PDF Documents: {pdf_count:,}\n"
                             f"• Office Documents: {office_count:,}\n"
                             f"• Other Files: {total_files - image_count - pdf_count - office_count:,}\n\n"
                             f"SEARCH TIPS:\n"
                             f"• Use 'Both' mode for comprehensive search\n"
                             f"• Use 'Filename Only' to find files by name\n"
                             f"• Use 'Content Only' to search inside documents\n"
                             f"• All file types are searchable by filename\n"
                             f"• Try searching: photo, document, image, etc.")
            
            # Create message box with explicit styling for better visibility
            msg_box = QMessageBox(self)
            msg_box.setWindowTitle("Scan Complete - DeepFileX")
            msg_box.setText(completion_msg)
            msg_box.setIcon(QMessageBox.Icon.Information)
            msg_box.setStandardButtons(QMessageBox.StandardButton.Ok)
            
            # Apply basic styling to ensure text is visible
            msg_box.setStyleSheet("""
                QMessageBox {
                    background-color: white;
                    color: black;
                    font-size: 12px;
                }
                QMessageBox QLabel {
                    color: black;
                    background-color: transparent;
                    font-size: 12px;
                }
                QMessageBox QPushButton {
                    background-color: #2c5aa0;
                    color: white;
                    border: none;
                    padding: 8px 16px;
                    border-radius: 5px;
                    font-weight: bold;
                    font-size: 12px;
                    min-width: 80px;
                }
                QMessageBox QPushButton:hover {
                    background-color: #1e3d72;
                }
            """)
            
            msg_box.exec()
            
            # 🆕 SmartLinks 컨텍스트 업데이트 (인덱싱 완료 후)
            if SMARTLINKS_AVAILABLE:
                try:
                    stats = self.search_engine.get_stats()
                    total_files = stats.get('total_files', count)
                    
                    if total_files > 50000:
                        self.update_smartlinks_context("system_optimization_needed")
                    else:
                        self.update_smartlinks_context("system_health_good")
                except Exception as e:
                    logger.error(f"SmartLinks 컨텍스트 업데이트 실패: {e}")
        
        except Exception as e:
            # Fallback message if there's any error building the detailed message
            logger.error(f"Error building completion message: {e}")
            fallback_msg = (f"Scan completed successfully!\n\n"
                           f"Files processed: {count:,}\n"
                           f"You can now search through your indexed files.")
            
            # Create fallback message box with explicit styling
            msg_box = QMessageBox(self)
            msg_box.setWindowTitle("Scan Complete")
            msg_box.setText(fallback_msg)
            msg_box.setIcon(QMessageBox.Icon.Information)
            msg_box.setStandardButtons(QMessageBox.StandardButton.Ok)
            
            # Apply basic styling to ensure text is visible
            msg_box.setStyleSheet("""
                QMessageBox {
                    background-color: white;
                    color: black;
                    font-size: 12px;
                }
                QMessageBox QLabel {
                    color: black;
                    background-color: transparent;
                    font-size: 12px;
                }
                QMessageBox QPushButton {
                    background-color: #2c5aa0;
                    color: white;
                    border: none;
                    padding: 8px 16px;
                    border-radius: 5px;
                    font-weight: bold;
                    font-size: 12px;
                    min-width: 80px;
                }
                QMessageBox QPushButton:hover {
                    background-color: #1e3d72;
                }
            """)
            
            msg_box.exec()
    
    def on_indexing_error(self, error_msg):
        """Handle indexing errors with enhanced progress bar handling"""
        self.index_btn.setText("Start Scan")
        
        # Show error in progress bar briefly
        self.progress_bar.setValue(0)
        self.progress_bar.setFormat("Error occurred - Indexing stopped")
        
        # Hide progress bar after showing error
        QTimer.singleShot(3000, lambda: self.progress_bar.setVisible(False))
        
        self.status_bar.showMessage("Indexing failed - Check error message")
        
        QMessageBox.critical(self, "Indexing Error", f"Indexing failed:\n{error_msg}")
    
    def clear_index(self):
        """Clear the search index"""
        reply = QMessageBox.question(self, "Clear Records", 
                                   "Are you sure you want to clear all scanned records?",
                                   QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
        
        if reply == QMessageBox.StandardButton.Yes:
            self.search_engine.clear_index()
            self.update_stats()
            self.results_list.clear()
            self.preview_text.clear()
            self.preview_label.setText("No file selected")
            self.status_bar.showMessage("Records cleared")
    
    def on_search_changed(self):
        """Handle search input changes"""
        # Debounced search
        self.search_timer.stop()
        if self.search_input.text().strip():
            self.search_timer.start(500)  # 500ms delay
    
    def perform_search(self):
        """Perform turbo search with mode selection"""
        query = self.search_input.text().strip()
        if not query:
            self.results_list.clear()
            self.results_header.setText("Diagnosis Results (0)")
            return
        
        # Get search mode
        mode_text = self.search_mode_combo.currentText()
        if "Content Only" in mode_text:
            search_mode = 'content'
        elif "Filename Only" in mode_text:
            search_mode = 'filename'
        else:
            search_mode = 'both'
        
        # Apply file type filter (except for filename-only search)
        file_type = self.file_type_combo.currentText()
        
        # Perform search
        results = self.search_engine.search_turbo(query, search_mode=search_mode, max_results=500)
        
        # Filter by file type if specified and not filename-only search
        if file_type != 'All Types' and search_mode != 'filename':
            results = [r for r in results if r['path'].lower().endswith(file_type.lower())]
        elif file_type != 'All Types' and search_mode == 'both':
            # For 'both' mode, only filter content matches, keep filename matches
            filtered_results = []
            for r in results:
                if r['path'].lower().endswith(file_type.lower()):
                    filtered_results.append(r)
                elif r.get('filename_match', False) and not r.get('content_match', False):
                    # Pure filename match - keep regardless of file type
                    filtered_results.append(r)
            results = filtered_results
        
        # Update results
        self.display_search_results(results, query)
    
    def display_search_results(self, results: List[Dict[str, Any]], query: str):
        """Display search results with match type indicators and file type icons"""
        self.results_list.clear()
        
        if not results:
            self.results_header.setText("Diagnosis Results (0)")
            return
        
        self.results_header.setText(f"Diagnosis Results ({len(results)})")
        
        for result in results:
            item = QListWidgetItem()
            
            # Format display text
            filename = result['filename']
            file_path = result['path']
            snippet = result['snippet']
            match_type = result.get('match_type', 'Content')
            
            # Get file type icon
            ext = Path(file_path).suffix.lower()
            if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg']:
                file_icon = "IMAGE"
            elif ext in ['.pdf']:
                file_icon = "PDF"
            elif ext in ['.docx', '.doc']:
                file_icon = "DOC"
            elif ext in ['.xlsx', '.xls']:
                file_icon = "XLS"
            elif ext in ['.pptx', '.ppt']:
                file_icon = "PPT"
            elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
                file_icon = "ZIP"
            else:
                file_icon = "FILE"
            
            # Truncate long paths
            if len(file_path) > 80:
                display_path = "..." + file_path[-77:]
            else:
                display_path = file_path
            
            # Truncate long snippets
            if len(snippet) > 150:
                snippet = snippet[:147] + "..."
            
            # Format with match type indicator and file icon
            item_text = f"{file_icon} {filename} [{match_type}]\nPath: {display_path}\nSnippet: {snippet}"
            item.setText(item_text)
            item.setData(Qt.ItemDataRole.UserRole, result)
            
            self.results_list.addItem(item)
        
        # 🆕 SmartLinks 컨텍스트 업데이트 (검색 결과에 따라)
        if SMARTLINKS_AVAILABLE:
            if len(results) > 1000:
                self.update_smartlinks_context("large_files_found")
            elif len(results) == 0:
                self.update_smartlinks_context("search_results_empty")
            else:
                self.update_smartlinks_context("file_scan_complete")
    
    def on_result_selected(self, item):
        """Handle result selection"""
        result = item.data(Qt.ItemDataRole.UserRole)
        if result:
            self.show_file_preview(result['path'])
    
    def show_file_preview(self, file_path: str):
        """Show file preview"""
        try:
            self.preview_label.setText(f"Preview: {Path(file_path).name}")
            
            # Extract content for preview
            extracted = self.extractor.extract_fast(file_path)
            if extracted['success']:
                content = extracted['content']
                if len(content) > 5000:
                    content = content[:5000] + "\n\n... (truncated for preview)"
                self.preview_text.setPlainText(content)
            else:
                self.preview_text.setPlainText("Could not extract content for preview.")
                
        except Exception as e:
            self.preview_text.setPlainText(f"Error loading preview: {str(e)}")
    
    def open_file_location(self, item):
        """Open file location in explorer"""
        result = item.data(Qt.ItemDataRole.UserRole)
        if result:
            file_path = result['path']
            try:
                if sys.platform == 'win32':
                    os.startfile(os.path.dirname(file_path))
                elif sys.platform == 'darwin':
                    os.system(f'open "{os.path.dirname(file_path)}"')
                else:
                    os.system(f'xdg-open "{os.path.dirname(file_path)}"')
            except Exception as e:
                QMessageBox.warning(self, "Error", f"Could not open location:\n{str(e)}")
    
    def clear_search(self):
        """Clear search input and results"""
        self.search_input.clear()
        self.results_list.clear()
        self.preview_text.clear()
        self.preview_label.setText("No file selected")
        self.results_header.setText("Diagnosis Results (0)")
    
    def update_stats(self):
        """Update statistics display"""
        stats = self.search_engine.get_stats()
        
        stats_text = f"""
Total Files: {stats['total_files']:,}
Content Words: {stats['total_words']:,}
Filename Words: {stats.get('total_filename_words', 0):,}
Last Update: {stats['last_update'] or 'Never'}
Search Paths: {len(self.search_paths)}
        """.strip()
        
        self.stats_label.setText(stats_text)
    
    def save_settings(self):
        """Save application settings"""
        self.settings.setValue('search_paths', self.search_paths)
        self.settings.setValue('geometry', self.saveGeometry())
        self.settings.setValue('window_state', self.saveState())
        self.settings.setValue('dark_mode', self.dark_mode)
    
    def load_settings(self):
        """Load application settings"""
        self.search_paths = self.settings.value('search_paths', [])
        self.update_folder_display()
        
        geometry = self.settings.value('geometry')
        if geometry:
            self.restoreGeometry(geometry)
        
        window_state = self.settings.value('window_state')
        if window_state:
            self.restoreState(window_state)
        
        # Load dark mode setting
        self.dark_mode = self.settings.value('dark_mode', False, type=bool)
        self.dark_mode_btn.setText("Light Mode" if self.dark_mode else "Dark Mode")
    
    def closeEvent(self, event):
        """Handle application close"""
        if self.indexing_worker and self.indexing_worker.isRunning():
            reply = QMessageBox.question(self, "Indexing in Progress", 
                                       "Indexing is in progress. Stop and exit?",
                                       QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No)
            
            if reply == QMessageBox.StandardButton.Yes:
                self.indexing_worker.stop()
                self.indexing_worker.wait()
            else:
                event.ignore()
                return
        
        self.save_settings()
        event.accept()
    
    # 🆕 SmartLinks 수익화 시스템 메서드들
    def init_smartlinks(self):
        """🎯 SmartLinks 시스템 초기화"""
        try:
            # SmartLinks 관리자 초기화
            self.smartlinks_manager = DeepFileXSmartLinksManager()
            
            # 광고 위젯 생성 (하단 배너)
            self.smartlinks_widget = AdBanner(
                parent=self,
                location="deepfilex_bottom_banner"
            )
            
            # 시그널 연결 (WebView2는 ad_clicked만 지원)
            # 람다로 감싸서 인자 개수 맞추기 (context 추가)
            self.smartlinks_widget.ad_clicked.connect(
                lambda url: self.on_smartlink_clicked("ad_banner", url)
            )
            
            # 메인 레이아웃에 추가 (제일 하단)
            main_layout = self.centralWidget().layout()
            main_layout.addWidget(self.smartlinks_widget)
            
            logger.info("🎯 SmartLinks 위젯 UI 통합 완료")
            
        except Exception as e:
            logger.error(f"❌ SmartLinks 초기화 실패: {e}")

    def on_smartlink_clicked(self, context, url):
        """💰 SmartLink 클릭 이벤트 처리"""
        try:
            # 로그 메시지
            logger.info(f"💰 광고 클릭: {context}")
            
            # 예상 수익 계산 및 표시
            if hasattr(self.smartlinks_manager, 'get_estimated_revenue'):
                revenue = self.smartlinks_manager.get_estimated_revenue()
                logger.info(f"📊 예상 누적 수익: ${revenue:.3f}")
            
            # 사용자에게 감사 메시지
            if hasattr(self, 'status_label'):
                self.status_label.setText("💝 DeepFileX 지원 감사합니다! 광고 수익으로 무료 서비스를 유지합니다.")
            
        except Exception as e:
            logger.error(f"SmartLink 클릭 처리 오류: {e}")

    def on_smartlink_shown(self, location):
        """👁️ 광고 노출 이벤트 처리"""
        logger.info(f"👁️ 광고 노출: {location}")

    def on_premium_requested(self):
        """👑 프리미엄 업그레이드 요청 처리"""
        logger.info("👑 프리미엄 업그레이드 요청됨")
        
        # 프리미엄 기능 활성화 후 광고 숨기기
        if hasattr(self, 'smartlinks_widget'):
            self.smartlinks_widget.hide()
            logger.info("✨ 프리미엄 모드 활성화 - 광고 제거됨")

    def update_smartlinks_context(self, new_context):
        """🎯 SmartLinks 컨텍스트 업데이트 (스캔 상태에 따라)"""
        if SMARTLINKS_AVAILABLE and hasattr(self, 'smartlinks_widget'):
            self.smartlinks_widget.set_context(new_context)
            logger.info(f"🔄 광고 컨텍스트 변경: {new_context}")

    def show_ad_settings(self):
        """광고 설정 다이얼로그 표시"""
        if hasattr(self, 'smartlinks_widget'):
            self.smartlinks_widget.show_ad_settings()

    def show_premium_upgrade(self):
        """프리미엄 업그레이드 다이얼로그 표시"""
        if hasattr(self, 'smartlinks_widget'):
            self.smartlinks_widget.show_premium_dialog()

    def show_revenue_stats(self):
        """수익 통계 표시"""
        if hasattr(self, 'smartlinks_manager'):
            stats = self.smartlinks_manager.session_stats
            revenue = self.smartlinks_manager.get_estimated_revenue()
            
            QMessageBox.information(self, "📊 SmartLinks 통계", 
                f"현재 세션 통계:\n\n"
                f"광고 노출: {stats['impressions']}회\n"
                f"광고 클릭: {stats['clicks']}회\n"
                f"클릭율: {self.smartlinks_manager.get_click_rate():.1f}%\n"
                f"예상 수익: ${revenue:.3f}\n\n"
                f"💝 DeepFileX를 지원해 주셔서 감사합니다!")

    # 🆕 자동 업데이트 시스템 메서드들
    def init_update_system(self):
        """🔄 자동 업데이트 시스템 초기화"""
        try:
            self.update_checker = UpdateChecker()
            
            # 신호 연결
            self.update_checker.update_available.connect(self.on_update_available)
            self.update_checker.no_update.connect(self.on_no_update)
            self.update_checker.error_occurred.connect(self.on_update_error)
            
            # 앱 시작 후 5초 뒤에 업데이트 체크 (UI 로딩 완료 후)
            startup_delay = UPDATE_CONFIG.get('startup_delay_seconds', 5) * 1000
            QTimer.singleShot(startup_delay, self.check_for_updates_startup)
            
            logger.info("🔄 업데이트 체커 초기화 완료")
            
        except Exception as e:
            logger.error(f"❌ 업데이트 시스템 초기화 실패: {e}")

    def check_for_updates_startup(self):
        """🔄 시작 시 업데이트 확인"""
        if hasattr(self, 'update_checker') and UPDATE_CONFIG.get('auto_check_enabled', True):
            logger.info("🔄 시작 시 업데이트 확인 중...")
            
            # 마지막 체크 시간 확인
            settings = QSettings('DeepFileX', 'Updates')
            last_check = settings.value('last_check_date')
            
            if last_check:
                from datetime import datetime, timedelta
                last_check_date = datetime.fromisoformat(last_check)
                check_interval = UPDATE_CONFIG.get('check_interval_days', 7)
                
                if datetime.now() - last_check_date < timedelta(days=check_interval):
                    logger.info(f"⏰ 다음 업데이트 체크까지 {check_interval}일 대기")
                    return
            
            self.update_checker.start()
    
    def on_update_available(self, update_info):
        """📦 업데이트 발견 시 처리"""
        logger.info(f"🎉 새로운 업데이트 발견: v{update_info.get('version', 'Unknown')}")
        
        try:
            # 업데이트 다이얼로그 표시
            dialog = UpdateDialog(self, update_info)
            dialog.exec()
            
            # 체크 완료 표시
            self.update_checker.mark_checked()
            
        except Exception as e:
            logger.error(f"❌ 업데이트 다이얼로그 표시 실패: {e}")
    
    def on_no_update(self):
        """✅ 업데이트 없음 처리"""
        logger.info("✅ 최신 버전 사용 중")
        
        # 체크 완료 표시
        if hasattr(self, 'update_checker'):
            self.update_checker.mark_checked()
    
    def on_update_error(self, error_message):
        """❌ 업데이트 체크 오류 처리"""
        logger.warning(f"⚠️ 업데이트 확인 실패: {error_message}")
        
        # 네트워크 연결 문제 등은 조용히 처리 (사용자에게 방해되지 않도록)

    def check_for_updates_manual(self):
        """🔄 수동 업데이트 확인 (메뉴에서 호출)"""
        if hasattr(self, 'update_manager'):
            # 강제 체크 (마지막 체크 시간 무시)
            self.update_manager.settings.setValue('last_check_date', '')
            self.update_manager.check_for_updates_async()
            
            # 사용자에게 확인 중임을 알림
            self.status_bar.showMessage("🔄 업데이트 확인 중...", 3000)
        else:
            QMessageBox.information(self, "업데이트 확인", 
                                   "업데이트 시스템이 초기화되지 않았습니다.")

    def show_update_settings(self):
        """🔄 업데이트 설정 다이얼로그"""
        if hasattr(self, 'update_manager'):
            from PyQt6.QtWidgets import QDialog, QVBoxLayout, QCheckBox, QSpinBox, QLabel, QFormLayout
            
            dialog = QDialog(self)
            dialog.setWindowTitle("🔄 업데이트 설정")
            dialog.setFixedSize(400, 250)
            
            layout = QVBoxLayout(dialog)
            form_layout = QFormLayout()
            
            settings = self.update_manager.settings
            
            # 자동 확인 설정
            auto_check = QCheckBox()
            auto_check.setChecked(settings.value('auto_check_enabled', True))
            form_layout.addRow("자동 업데이트 확인:", auto_check)
            
            # 확인 주기 설정
            check_interval = QSpinBox()
            check_interval.setRange(1, 30)
            check_interval.setValue(settings.value('check_interval_days', 7))
            check_interval.setSuffix(" 일")
            form_layout.addRow("확인 주기:", check_interval)
            
            # 백그라운드 확인 설정
            background_check = QCheckBox()
            background_check.setChecked(settings.value('background_check', True))
            form_layout.addRow("백그라운드 확인:", background_check)
            
            layout.addLayout(form_layout)
            
            # 버튼
            buttons = QHBoxLayout()
            ok_btn = QPushButton("확인")
            cancel_btn = QPushButton("취소")
            
            def save_settings():
                settings.setValue('auto_check_enabled', auto_check.isChecked())
                settings.setValue('check_interval_days', check_interval.value())
                settings.setValue('background_check', background_check.isChecked())
                dialog.accept()
            
            ok_btn.clicked.connect(save_settings)
            cancel_btn.clicked.connect(dialog.reject)
            
            buttons.addWidget(ok_btn)
            buttons.addWidget(cancel_btn)
            layout.addLayout(buttons)
            
            dialog.exec()



def main():
    """Main application entry point"""
    if not PYQT_AVAILABLE:
        print("PyQt6 is required but not installed.")
        print("Please install it with: pip install PyQt6")
        return
    
    app = QApplication(sys.argv)
    app.setApplicationName("DeepFileX")
    app.setOrganizationName("QuantumLayer")
    
    # Set application icon (if available)
    try:
        app.setWindowIcon(QIcon(""))
    except:
        pass
    
    # Create and show main window
    window = DeepFileX()
    window.show()
    
    # Check for PDF support
    if not PDF_AVAILABLE:
        QMessageBox.information(window, "PDF Support", 
                              "PDF support not available. To enable PDF content extraction, install:\n\n"
                              "Option 1: pip install PyPDF2\n"
                              "Option 2: pip install pdfplumber\n\n"
                              "Note: PDF files will still be indexed by filename without these libraries.")
    else:
        logger.info(f"PDF support enabled using {PDF_METHOD}")
        window.status_bar.showMessage(f"PDF support: {PDF_METHOD} available", 3000)
    
    # Check for Office document support
    missing_libs = []
    if not DOCX_AVAILABLE:
        missing_libs.append("python-docx (for Word documents)")
    if not EXCEL_AVAILABLE:
        missing_libs.append("openpyxl (for Excel files)")
    if not PPTX_AVAILABLE:
        missing_libs.append("python-pptx (for PowerPoint files)")
    
    if missing_libs:
        libs_info = "\n".join([f"pip install {lib.split(' ')[0]}" for lib in missing_libs])
        missing_info = "\n".join([f"• {lib}" for lib in missing_libs])
        
        QMessageBox.information(window, "Office Documents Support", 
                              f"For complete office document support, install:\n\n{libs_info}\n\n" +
                              f"Missing libraries:\n{missing_info}\n\n" +
                              "Note: Files will still be indexed by filename even without these libraries.")
    
    logger.info(f"DeepFileX started successfully!")
    logger.info(f"Library support - PDF: {PDF_AVAILABLE} ({PDF_METHOD if PDF_AVAILABLE else 'None'})")
    logger.info(f"Library support - DOCX: {DOCX_AVAILABLE}, EXCEL: {EXCEL_AVAILABLE}, PPTX: {PPTX_AVAILABLE}")
    logger.info(f"Supported image extensions: .jpg, .jpeg, .png, .gif, .bmp, .tiff, .webp, .ico, .svg")
    logger.info(f"Image files will be indexed by filename for fast searching")
    
    # Test image file detection
    test_image_exts = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.ico', '.svg']
    logger.info(f"Testing image file support: {', '.join(test_image_exts)}")
    
    sys.exit(app.exec())


if __name__ == "__main__":
    main()
